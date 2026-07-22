"""
Fidelity regression suite for the conversion engine.

Runs with one command (or run-tests.bat) and fails if any placement invariant
regresses on the committed fixtures. Covers the text-PLACEMENT paths, where
alignment/overlap regressions live:

  * PDF -> PPTX          (positioned blocks + native tables)
  * PDF -> edited PDF    (editor export: positioned blocks, byte-identical
                          unedited pages)

plus lightweight structural smoke checks for the other from-PDF paths
(PDF -> images, PDF -> text). Office-COM paths are covered by test_com_crash /
verify_outputs and are intentionally not re-run here (Office layout has no
text-placement fidelity for us to assert, and the suite must run without Office).

Invariants asserted per PPTX slide:
  1. Paragraph-level blocks — no sub-0.3in text box that continues an adjacent
     paragraph (fragmentation).
  2. No fabricated ' / ' line-join artifacts.
  3. Every table the PDF detects becomes a NATIVE PowerPoint table.
  4. Zero text-box insets (all four margins 0) so positioning is baseline-true.
  5. Consistent font mapping — each source font maps to exactly one target.
  6. No two elements overlap by >10% of the smaller element's area.
  7. No element extends past the page boundary.
  8. Text-box positions stay within tolerance of the committed golden layout.

Usage:
  python tests/fidelity_suite.py                 # run the suite
  python tests/fidelity_suite.py --update-golden # re-bless golden layout
"""

import glob
import io
import json
import os
import re
import sys

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, os.path.join(HERE, ".."))

import fitz
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.converter import (convert_pdf_to_pptx, _table_cells_tabular,
                           PROSE_CELL_CHARS, PROSE_MAX_TEXT_CELLS,
                           MIN_TABULAR_CELLS)
from app.edit_model import EditSession
from app.extraction import (_collect_lines, _infer_aligned_tables,
                            _inherit_glyph_colors, _point_in, _center,
                            _font_is_dingbat)
from app.pdf_edit_export import export_edited_pdf
from app.convert_pdf_misc import pdf_to_images, pdf_to_text
from tests.engines import ENGINES, EngineUnavailable

NL = chr(10)
FIX = os.path.join(HERE, "fixtures")
GOLD = os.path.join(HERE, "fidelity", "golden.json")
# Writable work area; overridable so a frozen/packaged build can point it at a
# per-user temp dir instead of the (read-only) app bundle.
WORK = os.environ.get("BENCHPDF_DIAG_WORK") or os.path.join(HERE, "..", "_work", "fidelity")
os.makedirs(WORK, exist_ok=True)

FRAG_MAX_H_IN = 0.30          # boxes shorter than this are fragmentation suspects
OVERLAP_FRAC = 0.10           # >10% of smaller element's area = overlap
BOUND_EPS_IN = 0.06           # allowed slop past a page edge
POS_TOL = 0.02                # golden position drift tolerance (fraction of page)


class Result:
    def __init__(self):
        self.checks = []      # (ok, label)
        self.skips = []

    def check(self, ok, label):
        self.checks.append((bool(ok), label))
        return ok

    def skip(self, label):
        self.skips.append(label)

    @property
    def passed(self):
        return all(ok for ok, _ in self.checks)


def _emu_in(v):
    return v / 914400.0


def slide_elements(slide):
    """Return [(kind, x, y, w, h, text)] for text boxes and native tables."""
    els = []
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            els.append(("text", sh.left, sh.top, sh.width, sh.height,
                        " ".join(sh.text_frame.text.split())[:24]))
        elif sh.has_table:
            els.append(("table", sh.left, sh.top, sh.width, sh.height, "TABLE"))
    return els


def overlap_frac(a, b):
    ix = max(0, min(a[1] + a[3], b[1] + b[3]) - max(a[1], b[1]))
    iy = max(0, min(a[2] + a[4], b[2] + b[4]) - max(a[2], b[2]))
    inter = ix * iy
    if inter <= 0:
        return 0.0
    return inter / max(1, min(a[3] * a[4], b[3] * b[4]))


def check_text_insets(prs, res, name):
    bad = 0
    for slide in prs.slides:
        for sh in slide.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                tf = sh.text_frame
                if any(m not in (0, None) and _emu_in(m or 0) > 0.001
                       for m in (tf.margin_left, tf.margin_right, tf.margin_top, tf.margin_bottom)):
                    bad += 1
    res.check(bad == 0, f"[{name}] text-box insets are zero (offenders: {bad})")


def _first_run_size(sh):
    try:
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if r.font.size is not None:
                    return r.font.size.pt
    except Exception:
        pass
    return None


def check_fragments(prs, res, name):
    """A real fragment is a wrapped line of a paragraph broken into its own box:
    short, sharing the previous box's left edge and font size, a tiny gap below
    it, and the upper text ending mid-sentence (no terminal punctuation). Short
    headings differ (larger/different font size than the body beneath them) and
    complete lines end in punctuation, so neither is flagged."""
    suspects = 0
    for slide in prs.slides:
        tb = []
        for sh in slide.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                x, y, w, h = sh.left, sh.top, sh.width, sh.height
                tb.append((x, y, w, h, " ".join(sh.text_frame.text.split()),
                           _first_run_size(sh)))
        for a in tb:
            if _emu_in(a[3]) >= FRAG_MAX_H_IN:
                continue
            a_ends_clean = a[4][-1:] in ".!?:;)’\"" if a[4] else True
            if a_ends_clean:
                continue
            for b in tb:
                if b is a:
                    continue
                same_left = abs(_emu_in(a[0] - b[0])) < 0.05
                gap = _emu_in(b[1] - (a[1] + a[3]))
                below = -0.05 <= gap <= 0.06
                same_size = (a[5] is not None and b[5] is not None
                             and abs(a[5] - b[5]) <= 0.6)
                if same_left and below and same_size:
                    suspects += 1
                    break
    res.check(suspects == 0, f"[{name}] no fragmented paragraphs (suspect boxes: {suspects})")


def check_slashes(prs, src_pdf, res, name):
    doc = fitz.open(src_pdf)
    src = "\n".join(doc[i].get_text("text") for i in range(doc.page_count))
    doc.close()
    out = "\n".join(sh.text_frame.text for slide in prs.slides for sh in slide.shapes
                    if sh.has_text_frame)
    res.check(out.count(" / ") <= src.count(" / "),
              f"[{name}] no fabricated ' / ' joins (out={out.count(' / ')}, src={src.count(' / ')})")


def _native_tables(prs):
    return sum(1 for sl in prs.slides for sh in sl.shapes if sh.has_table)


def _source_tables(src_pdf):
    """Expected tables per page, engine-independently: ruled tables that
    CONTAIN text (a text-less decorative grid is graphics, not a table),
    plus unruled tables recovered by column-alignment inference, minus every
    grid whose cells hold prose rather than tabular content (a bordered
    callout is furniture however it was detected)."""
    doc = fitz.open(src_pdf)
    per_page = []
    for i in range(doc.page_count):
        lines = _collect_lines(doc[i].get_text("dict"))
        try:
            found = doc[i].find_tables(strategy="lines").tables
        except Exception:
            found = []
        ruled = [t for t in found
                 if t.row_count >= 1 and t.col_count >= 1
                 and any(_point_in(t.bbox, *_center(ln["bbox"])) for ln in lines)]
        inferred = _infer_aligned_tables(lines, [t.bbox for t in ruled])
        ruled = [t for t in ruled if _table_cells_tabular(t, lines)]
        inferred = [t for t in inferred if _table_cells_tabular(t, lines)]
        per_page.append({"ruled": ruled, "inferred": inferred, "lines": lines})
    doc.close()
    return per_page


def check_tables(prs, src_pdf, res, name):
    """Every table the PDF detects must arrive as a NATIVE PowerPoint table -
    including UNRULED tables recovered from column alignment (the statement
    ledger class). Derived from the .pptx and the source PDF, not from any
    engine's own report, so the same assertion holds for an engine that
    reports nothing.
    """
    pages = _source_tables(src_pdf)
    want = sum(len(p["ruled"]) + len(p["inferred"]) for p in pages)
    got = _native_tables(prs)
    res.check(got >= want,
              f"[{name}] every PDF table is native, ruled or not (source {want}, output {got})")

    # each inferred (unruled) source table's text must live in a NATIVE table,
    # not in loose text boxes
    table_text = " ".join(
        " ".join(c.text.split())
        for sl in prs.slides for sh in sl.shapes if sh.has_table
        for row in sh.table.rows for c in row.cells)
    missing = []
    for p in pages:
        for t in p["inferred"]:
            probe = next((" ".join("".join(s["text"] for s in ln["spans"]).split())
                          for ln in p["lines"]
                          if _point_in(t.bbox, *_center(ln["bbox"]))
                          and len("".join(s["text"] for s in ln["spans"]).strip()) >= 6),
                         None)
            if probe and probe not in table_text:
                missing.append(probe[:30])
    if any(p["inferred"] for p in pages):
        res.check(not missing,
                  f"[{name}] unruled ledger text lands in native tables (missing: {missing[:3]})")


def check_no_empty_tables(prs, res, name):
    """A native table with no text in any cell is a phantom: a decorative
    grid that should have stayed graphics (fault: 8-square brand mark
    emitted as an empty 1x8 table)."""
    empty = 0
    for sl in prs.slides:
        for sh in sl.shapes:
            if sh.has_table:
                if not any(c.text.strip() for row in sh.table.rows for c in row.cells):
                    empty += 1
    res.check(empty == 0, f"[{name}] no native table with all-empty cells (found: {empty})")


def check_no_prose_tables(prs, res, name):
    """A bordered callout must not ship as a native table.

    Read off the OUTPUT alone, so it holds for any engine: a shipped table
    whose cells are nearly all empty, or which is one paragraph plus a label,
    is a box drawn around prose. Inside a table that prose reflows into a
    column one cell wide (fault: a ruled 2x2 around a whole guidance column,
    the real corpus's worst page at 0.1414).
    """
    bad = []
    for si, sl in enumerate(prs.slides, 1):
        for sh in sl.shapes:
            if not sh.has_table:
                continue
            lens = [len(c.text.strip()) for row in sh.table.rows for c in row.cells]
            filled = [n for n in lens if n > 0]
            if len(filled) < MIN_TABULAR_CELLS:
                bad.append(f"slide {si}: {len(filled)} cell(s) with text")
            elif len(filled) <= PROSE_MAX_TEXT_CELLS and max(lens) >= PROSE_CELL_CHARS:
                bad.append(f"slide {si}: {max(lens)} chars in one of {len(filled)} filled cells")
    res.check(not bad, f"[{name}] no bordered callout shipped as a table ({bad[:2]})")


def check_lists(prs, src_pdf, res, name):
    """A bulleted list must arrive as one paragraph per item, and its marker
    must arrive as a bullet.

    The W3C WCAG working draft sets every marker as a 7pt ZapfDingbats "H",
    which draws a hollow circle. Nothing about the code point says "bullet",
    so nothing flagged the lines as list items and 25 entries clustered into
    ONE reflowing paragraph - the real corpus's worst page. Counted off the
    source PDF and the .pptx, so it holds for any engine.
    """
    # The source count is computed HERE, from font names and code points in the
    # PDF, and deliberately does NOT call the engine's marker detector: routing
    # it through _attach_markers would make the expectation collapse to zero
    # whenever the detector broke, and the check would skip instead of fail.
    dingbat_tokens = ("zapfdingbats", "dingbat", "wingding", "webding")
    bullets = "•◦‣·∙▪▫◾"
    doc = fitz.open(src_pdf)
    want = 0
    for i in range(doc.page_count):
        for blk in doc[i].get_text("dict")["blocks"]:
            for ln in blk.get("lines", []):
                spans = [s for s in ln.get("spans", []) if (s.get("text") or "").strip()]
                if len(spans) < 2:
                    continue
                glyph = spans[0]["text"].strip()
                if len(glyph) != 1:
                    continue
                fname = (spans[0].get("font") or "").lower()
                if (any(t in fname for t in dingbat_tokens)
                        or glyph in bullets or 0xF000 <= ord(glyph) <= 0xF0FF):
                    want += 1
    doc.close()
    if not want:
        res.skip(f"[{name}] no list markers in source")
        return

    got, dingbats = 0, set()
    for sl in prs.slides:
        for sh in sl.shapes:
            frames = []
            if sh.has_text_frame:
                frames.append(sh.text_frame)
            elif sh.has_table:
                frames.extend(c.text_frame for row in sh.table.rows for c in row.cells)
            for tf in frames:
                for para in tf.paragraphs:
                    if para.text.lstrip().startswith("•"):
                        got += 1
                    for r in para.runs:
                        if _font_is_dingbat(r.font.name):
                            dingbats.add(r.font.name)
    res.check(got >= want,
              f"[{name}] every list item keeps its own paragraph (source {want}, output {got})")
    # a dingbat glyph mapped into a text font is a stray letter beside the item
    res.check(not dingbats,
              f"[{name}] list markers normalised, not shipped as dingbats ({sorted(dingbats)})")


def _editable_text(prs):
    parts = []
    for sl in prs.slides:
        for sh in sl.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                parts.append(" ".join(sh.text_frame.text.split()))
            elif sh.has_table:
                for row in sh.table.rows:
                    for c in row.cells:
                        if c.text.strip():
                            parts.append(" ".join(c.text.split()))
    return " ".join(parts)


def check_no_ghost_text(prs, src_pdf, res, name):
    """No text may exist BOTH as pixels in a background raster and as editable
    text on top: that is the ghost-doubling fault (form-XObject/outline-path
    statements). OCR-free: for every source line whose text is present in the
    editable layer, the full-page background image under its bbox must be
    (near) ink-free. Fixture backgrounds are plain under text, so ink there is
    the line's own ghost."""
    from PIL import Image
    editable = _editable_text(prs)
    doc = fitz.open(src_pdf)
    sw_pt = prs.slide_width / 12700.0
    offenders = 0
    checked = 0
    for i, slide in enumerate(prs.slides):
        if i >= doc.page_count:
            break
        # the background is a picture covering (nearly) the whole slide
        bg = None
        for sh in slide.shapes:
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE \
                    and sh.width >= 0.9 * prs.slide_width \
                    and sh.height >= 0.9 * prs.slide_height:
                bg = sh
        if bg is None:
            continue
        img = Image.open(io.BytesIO(bg.image.blob)).convert("RGB")
        z = img.width / sw_pt
        for ln in _collect_lines(doc[i].get_text("dict")):
            txt = " ".join("".join(s["text"] for s in ln["spans"]).split())
            if len(txt) < 6 or txt not in editable:
                continue
            x0, y0, x1, y1 = ln["bbox"]
            px0, py0 = max(int(x0 * z), 0), max(int(y0 * z), 0)
            px1 = min(int(x1 * z), img.width)
            py1 = min(int(y1 * z), img.height)
            if px1 - px0 < 2 or py1 - py0 < 2:
                continue
            crop = img.crop((px0, py0, px1, py1))
            px = list(crop.getdata())
            dark = sum(1 for p in px if p[0] < 128 and p[1] < 128 and p[2] < 128)
            checked += 1
            if dark / len(px) > 0.05:
                offenders += 1
    doc.close()
    if checked:
        res.check(offenders == 0,
                  f"[{name}] no text doubled between background raster and editable layer "
                  f"(ghosted lines: {offenders}/{checked})")
    else:
        res.skip(f"[{name}] no full-page background raster to check for ghosts")


_MEASURE_FONTS = {}


def _measure_word(word, size, bold, mono):
    key = ("cour" if mono else ("hebo" if bold else "helv"))
    if key not in _MEASURE_FONTS:
        _MEASURE_FONTS[key] = fitz.Font(key)
    return _MEASURE_FONTS[key].text_length(word, fontsize=size)


def check_no_midword_wrap(prs, res, name):
    """A wrapping text box must be wide enough for its longest word, or
    PowerPoint breaks the word across lines (fault: narrow header blocks
    like 'END OF DAY ACCOUNT BALANCE' wrapping mid-word). Word widths are
    estimated with metric-compatible base-14 fonts."""
    offenders = []
    for slide in prs.slides:
        for sh in slide.shapes:
            if not (sh.has_text_frame and sh.text_frame.text.strip()):
                continue
            tf = sh.text_frame
            if tf.word_wrap is False:      # wrap off: PPT never breaks words
                continue
            w_pt = sh.width / 12700.0
            for para in tf.paragraphs:
                for r in para.runs:
                    size = r.font.size.pt if r.font.size else 10.0
                    bold = bool(r.font.bold)
                    fname = (r.font.name or "").lower()
                    mono = any(k in fname for k in ("consol", "courier", "mono"))
                    for word in r.text.split():
                        if _measure_word(word, size, bold, mono) > w_pt + 0.75:
                            offenders.append(word[:16])
    res.check(not offenders,
              f"[{name}] no box narrower than its longest word "
              f"(mid-word wrap risks: {offenders[:4]})")


def _source_line_colors(src_pdf):
    """normalised line text -> rendered ink colour, glyph-outline aware."""
    doc = fitz.open(src_pdf)
    out = {}
    ambiguous = set()
    for i in range(doc.page_count):
        page = doc[i]
        lines = _collect_lines(page.get_text("dict"))
        try:
            drawings = page.get_drawings()
        except Exception:
            drawings = []
        _inherit_glyph_colors(lines, drawings, page.rect.width, page.rect.height)
        for ln in lines:
            txt = " ".join("".join(s["text"] for s in ln["spans"]).split())
            if len(txt) < 6:
                continue
            spans = [s for s in ln["spans"] if s["text"].strip()]
            cols = {int(s.get("color", 0)) & 0xFFFFFF for s in spans}
            if len(cols) != 1:
                continue          # mixed-colour lines can't be asserted at line level
            col = cols.pop()
            if txt in out and out[txt] != col:
                ambiguous.add(txt)
            out[txt] = col
    doc.close()
    return {t: c for t, c in out.items() if t not in ambiguous}


def _run_colors(prs):
    """[(normalised paragraph text, run colour int)] for text boxes and cells."""
    out = []

    def eat(tf):
        for para in tf.paragraphs:
            ptxt = " ".join("".join(r.text for r in para.runs).split())
            for r in para.runs:
                try:
                    col = r.font.color.rgb          # RGBColor is a str subclass
                    col = int(str(col), 16)
                except Exception:
                    col = 0
                out.append((ptxt, col))

    for sl in prs.slides:
        for sh in sl.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                eat(sh.text_frame)
            elif sh.has_table:
                for row in sh.table.rows:
                    for c in row.cells:
                        if c.text.strip():
                            eat(c.text_frame)
    return out


def check_colors(prs, src_pdf, res, name):
    """Output run colours must match the SOURCE'S RENDERED ink colour within
    tolerance: greys and brand colours survive, nothing flattens to #000000.
    Source truth is glyph-outline aware, so invisible-text-layer statements
    are held to the colour of the ink actually painted."""
    src = _source_line_colors(src_pdf)
    runs = _run_colors(prs)
    mismatches = []
    matched_src_colors, matched_out_colors = set(), set()
    for txt, want in src.items():
        hits = [c for t, c in runs if t == txt]
        if not hits:
            continue
        wr, wg, wb = (want >> 16) & 255, (want >> 8) & 255, want & 255
        ok = any(abs(((c >> 16) & 255) - wr) <= 40
                 and abs(((c >> 8) & 255) - wg) <= 40
                 and abs((c & 255) - wb) <= 40 for c in hits)
        matched_src_colors.add(want)
        if ok:
            matched_out_colors.update(hits)
        else:
            mismatches.append((txt[:24], "#%06x" % want, "#%06x" % hits[0]))
    res.check(not mismatches,
              f"[{name}] run colours match source ink within tolerance "
              f"(mismatched: {mismatches[:3]})")
    if len(matched_src_colors) >= 2:
        res.check(len(set(matched_out_colors)) >= 2,
                  f"[{name}] palette not flattened to one colour "
                  f"(source {len(matched_src_colors)} colours, output {len(set(matched_out_colors))})")


def _family(font_name):
    """Reduce a PDF font name to its family.

    PDF font names encode weight and style ('Helvetica-Bold', 'ABCDEF+Arial,BoldItalic')
    while PowerPoint carries those as run properties, so comparing raw names would
    call a correct conversion a collapse. Compare families; check weight separately.
    """
    n = re.sub(r"^[A-Z]{6}[+]", "", font_name or "")
    n = re.split(r"[-,]", n)[0]
    return re.sub(r"(?i)(bold|italic|oblique|light|regular|medium|semibold)$", "", n).strip().lower()


def check_fonts(prs, src_pdf, res, name):
    """The output must not collapse every source family onto one font.

    A converter that hardcodes a single fontFace passes every positional check
    while destroying the document's typography. Asserted against the .pptx and
    the source PDF, so it holds for any engine, reporting or not.
    """
    doc = fitz.open(src_pdf)
    src = set()
    for i in range(doc.page_count):
        for blk in doc[i].get_text("dict")["blocks"]:
            for ln in blk.get("lines", []):
                for sp in ln["spans"]:
                    # A dingbat font is not typography, it is list markers, and
                    # the engine deliberately normalises those to a bullet in
                    # the text's own font. Counting it as a source family would
                    # make correct normalisation read as a collapse.
                    if _font_is_dingbat(sp["font"]):
                        continue
                    src.add(_family(sp["font"]))
    doc.close()

    out, bold_runs = set(), 0
    for sl in prs.slides:
        for sh in sl.shapes:
            if not sh.has_text_frame:
                continue
            for para in sh.text_frame.paragraphs:
                for r in para.runs:
                    if r.font.name:
                        out.add(_family(r.font.name))
                    if r.font.bold:
                        bold_runs += 1

    want = min(len(src), 2)
    res.check(len(out) >= want,
              f"[{name}] font families not collapsed "
              f"(source {len(src)}, output {len(out)}: {sorted(out)})")

    # Weight is the other half of typography, and the half a family-level check
    # would otherwise let through.
    doc = fitz.open(src_pdf)
    src_bold = any("bold" in (sp["font"] or "").lower()
                   for i in range(doc.page_count)
                   for blk in doc[i].get_text("dict")["blocks"]
                   for ln in blk.get("lines", [])
                   for sp in ln["spans"])
    doc.close()
    if src_bold:
        res.check(bold_runs > 0,
                  f"[{name}] bold weight survives (bold runs in output: {bold_runs})")
    else:
        res.skip(f"[{name}] no bold in source")


def check_graphics(prs, src_pdf, res, name):
    """Pages with a graphic layer must not arrive as bare text.

    The failure this catches: an engine that emits only text boxes, so fills,
    rules, charts and table shading silently vanish from every slide.
    """
    doc = fitz.open(src_pdf)
    pages_with_art = []
    for i in range(doc.page_count):
        page = doc[i]
        if len(page.get_drawings()) >= 3 or page.get_images():
            pages_with_art.append(i)
    doc.close()
    if not pages_with_art:
        res.skip(f"[{name}] no graphic layer in source")
        return
    carriers = sum(1 for sl in prs.slides for sh in sl.shapes
                   if sh.has_table or sh.shape_type == MSO_SHAPE_TYPE.PICTURE
                   or getattr(sh, "fill", None) is not None and not sh.has_text_frame)
    res.check(carriers > 0,
              f"[{name}] graphic layer preserved "
              f"({len(pages_with_art)} source pages have art, output carriers: {carriers})")


def check_degenerate_geometry(prs, res, name):
    """No zero-height text elements; no pile-ups at one origin.

    Guards the form-XObject bug class: an un-composed transform stack stacked
    footer text boxes at the page origin with zero height (fixture:
    form_xobject_statement).
    """
    zero_dim, piled = 0, 0
    for slide in prs.slides:
        origins = {}
        for sh in slide.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                if sh.height == 0 or sh.width == 0:
                    zero_dim += 1
                key = (sh.left, sh.top)
                origins[key] = origins.get(key, 0) + 1
        piled += sum(1 for v in origins.values() if v > 2)
    res.check(zero_dim == 0, f"[{name}] no zero-height or zero-width text elements (found: {zero_dim})")
    res.check(piled == 0, f"[{name}] no more than 2 text elements share an origin (piles: {piled})")
    return zero_dim == 0


def check_effective_sizes(prs, src_pdf, res, name):
    """Every output font size matches a SOURCE span's effective size within
    0.5pt, and nothing dips under 4pt unless the source truly has it. A
    transform-scale corruption fails both instantly."""
    doc = fitz.open(src_pdf)
    src_sizes = set()
    for i in range(doc.page_count):
        for blk in doc[i].get_text("dict")["blocks"]:
            for ln in blk.get("lines", []):
                for sp in ln["spans"]:
                    if sp["text"].strip():
                        src_sizes.add(round(sp["size"], 2))
    doc.close()

    bad, tiny = [], []
    for slide in prs.slides:
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            for para in sh.text_frame.paragraphs:
                for r in para.runs:
                    if r.font.size is None or not r.text.strip():
                        continue
                    pt = r.font.size.pt
                    if not any(abs(pt - ss) <= 0.5 for ss in src_sizes):
                        bad.append(round(pt, 2))
                    if pt < 4 and not any(ss < 4.5 for ss in src_sizes):
                        tiny.append(round(pt, 2))
    res.check(not bad, f"[{name}] output sizes match source within 0.5pt (alien: {sorted(set(bad))[:6]})")
    res.check(not tiny, f"[{name}] no sub-4pt sizes the source lacks (found: {sorted(set(tiny))[:6]})")


def check_overlap_bounds(prs, res, name):
    SW, SH = prs.slide_width, prs.slide_height
    overlaps, oob = 0, 0
    for si, slide in enumerate(prs.slides):
        els = slide_elements(slide)
        for i in range(len(els)):
            e = els[i]
            if (_emu_in(e[1]) < -BOUND_EPS_IN or _emu_in(e[2]) < -BOUND_EPS_IN
                    or _emu_in(e[1] + e[3] - SW) > BOUND_EPS_IN
                    or _emu_in(e[2] + e[4] - SH) > BOUND_EPS_IN):
                oob += 1
            for j in range(i + 1, len(els)):
                if overlap_frac(e, els[j]) > OVERLAP_FRAC:
                    overlaps += 1
    res.check(overlaps == 0, f"[{name}] no element overlaps >10% (found: {overlaps})")
    res.check(oob == 0, f"[{name}] no element past the page boundary (found: {oob})")


def layout_signature(prs):
    """Normalised text-box positions per slide for golden comparison."""
    SW, SH = prs.slide_width, prs.slide_height
    sig = []
    for slide in prs.slides:
        boxes = [(e[1] / SW, e[2] / SH, e[3] / SW, e[4] / SH)
                 for e in slide_elements(slide) if e[0] == "text"]
        boxes.sort()
        sig.append(boxes)
    return sig


def check_golden(prs, golden, res, name):
    if name not in golden:
        res.skip(f"[{name}] no golden layout yet")
        return
    cur = layout_signature(prs)
    ref = golden[name]
    if len(cur) != len(ref):
        res.check(False, f"[{name}] slide count matches golden ({len(cur)} vs {len(ref)})")
        return
    worst = 0.0
    countmis = 0
    for cs, rs in zip(cur, ref):
        if len(cs) != len(rs):
            countmis += 1
            continue
        for cb, rb in zip(cs, rs):
            # position AND size: a width error (lost or phantom advance) is
            # exactly how the form-XObject bug class shows up in output
            worst = max(worst, *(abs(a - b) for a, b in zip(cb, rb)))
    res.check(countmis == 0, f"[{name}] golden box counts per slide match (mismatched slides: {countmis})")
    res.check(worst <= POS_TOL, f"[{name}] text-box drift within {POS_TOL*100:.0f}% (worst: {worst*100:.2f}%)")


# --------------------------------------------------------------------------- #
def run_pptx(engine, update_golden=False):
    """Run the placement invariants for one engine over every fixture.

    The golden layout is the Python engine's blessed output. Holding a second
    engine to the same golden is exactly the parity assertion we want: if the
    browser engine emits 57 boxes where Python emits 5, the box-count check
    fails, which is the whole point.
    """
    print(NL + "== PDF -> PPTX [%s: %s] ==" % (engine.name, engine.label))
    res = Result()
    res.engine = engine.name
    golden = {}
    if os.path.exists(GOLD):
        golden = json.load(open(GOLD))
    new_golden = dict(golden)
    for src in sorted(glob.glob(os.path.join(FIX, "*.pdf"))):
        name = os.path.splitext(os.path.basename(src))[0]
        out = os.path.join(WORK, f"{name}.{engine.name}.pptx")
        try:
            engine.convert(src, out)
        except EngineUnavailable as e:
            res.skip(f"[{name}] engine '{engine.name}' unavailable: {e}")
            continue
        except Exception as e:
            res.check(False, f"[{name}] engine '{engine.name}' raised: {type(e).__name__}: {e}")
            continue
        prs = Presentation(out)
        geometry_sound = check_degenerate_geometry(prs, res, name)
        check_fragments(prs, res, name)
        check_slashes(prs, src, res, name)
        check_tables(prs, src, res, name)
        check_no_empty_tables(prs, res, name)
        check_no_prose_tables(prs, res, name)
        check_lists(prs, src, res, name)
        check_no_ghost_text(prs, src, res, name)
        check_no_midword_wrap(prs, res, name)
        check_colors(prs, src, res, name)
        check_text_insets(prs, res, name)
        check_fonts(prs, src, res, name)
        check_graphics(prs, src, res, name)
        check_effective_sizes(prs, src, res, name)
        if geometry_sound:
            check_overlap_bounds(prs, res, name)
        else:
            # overlap over zero-sized elements is vacuous; fail it loudly so a
            # broken geometry invariant can never mask a broken overlap one
            res.check(False, f"[{name}] overlap check NOT RUN: masked by degenerate geometry")
        if update_golden:
            new_golden[name] = layout_signature(prs)
        else:
            check_golden(prs, golden, res, name)
    if update_golden:
        os.makedirs(os.path.dirname(GOLD), exist_ok=True)
        json.dump(new_golden, open(GOLD, "w"), indent=1)
        print("  golden layout written to", os.path.relpath(GOLD, HERE))
    return res


def run_editor():
    print("\n== PDF -> edited PDF (editor export) ==")
    res = Result()
    for src in sorted(glob.glob(os.path.join(FIX, "*.pdf"))):
        name = os.path.splitext(os.path.basename(src))[0]
        # 1) no-edit export must be byte-identical everywhere
        out0 = os.path.join(WORK, name + "-noedit.pdf")
        export_edited_pdf(src, {"pages": {}}, out0)
        a = open(src, "rb").read()
        res.check(open(out0, "rb").read()[:len(a)] == a,
                  f"[{name}] no-edit export keeps original bytes")

        # 2) one edited page: edit persists, other pages byte-identical, clean text
        s = EditSession(src)
        model = s.page_model(0)
        blocks = model["blocks"]
        tgt = next((b for b in blocks if b["type"] == "text" and len(b["html"]) > 20), None)
        s.close()
        if tgt is None:
            res.skip(f"[{name}] page 1 has no text block to edit")
            continue
        end = tgt["html"].index("</div>") + 6 if "</div>" in tgt["html"] else len(tgt["html"])
        tgt["html"] = "<div>FIDELITY-EDIT sentinel value.</div>" + tgt["html"][end:]
        out1 = os.path.join(WORK, name + "-edited.pdf")
        export_edited_pdf(src, {"pages": {"0": {"blocks": blocks}}}, out1)
        o1, o2 = fitz.open(src), fitz.open(out1)
        ident = all(o1[i].read_contents() == o2[i].read_contents()
                    for i in range(o1.page_count) if i != 0)
        res.check(ident, f"[{name}] editor: unedited pages byte-identical")
        t = o2[0].get_text("text")
        res.check("FIDELITY-EDIT sentinel value" in " ".join(t.split()),
                  f"[{name}] editor: edited text present (normal spaces)")
        res.check(chr(0xA0) not in t and chr(0xAD) not in t,
                  f"[{name}] editor: no nbsp/soft-hyphen artifacts")
        o1.close(); o2.close()
    return res


def run_other_paths():
    """Structural smoke for the remaining conversion-hub paths that don't have
    text-placement fidelity to assert. Office-COM paths (Word/Excel/PPTX -> PDF,
    PDF -> Word) and the URL path are exercised by their own tests
    (tests/verify_outputs.py, tests/verify_hub.py, tests/test_com_crash.py) and
    are intentionally not re-run here so the suite stays fast and needs no Office
    or network. Those paths don't import the extraction/placement code changed
    in this area, so they can't be affected by a placement regression."""
    print("\n== other hub paths (structural) ==")
    res = Result()
    src = os.path.join(FIX, "text_report.pdf")

    imgs = pdf_to_images(src, os.path.join(WORK, "imgs"), fmt="png", dpi=120)
    doc = fitz.open(src); npages = doc.page_count; doc.close()
    res.check(len(imgs) == npages, f"PDF->images: one image per page ({len(imgs)}/{npages})")

    txt = os.path.join(WORK, "text_report.txt")
    pdf_to_text(src, txt)
    body = open(txt, encoding="utf-8").read()
    res.check("Quarterly Operations Review" in body, "PDF->text: extracted expected content")

    # images -> merged PDF (the TO-PDF image path; Pillow/img2pdf, no Office)
    imgdir = os.path.join(HERE, "input")
    pics = [os.path.join(imgdir, f) for f in ("photo_a.jpg", "photo_b.png")
            if os.path.exists(os.path.join(imgdir, f))]
    if len(pics) >= 2:
        from app.convert_images import images_to_pdf
        merged = os.path.join(WORK, "merged.pdf")
        images_to_pdf(pics, merged)
        d = fitz.open(merged); ok = d.page_count == len(pics); d.close()
        res.check(ok, f"images->PDF: merged {len(pics)} images into {len(pics)}-page PDF")
    else:
        res.skip("images->PDF (sample images not present)")

    res.skip("Office-COM paths (Word/Excel/PPTX->PDF, PDF->Word) — see verify_outputs.py / test_com_crash.py")
    res.skip("Web page->PDF — see verify_hub.py (needs network)")
    return res


def main():
    update = "--update-golden" in sys.argv

    # --engines python,browser   (default: every registered engine)
    sel = None
    for a in sys.argv:
        if a.startswith("--engines="):
            sel = [x.strip() for x in a.split("=", 1)[1].split(",") if x.strip()]
    names = sel or list(ENGINES)
    unknown = [n for n in names if n not in ENGINES]
    if unknown:
        print("Unknown engine(s):", ", ".join(unknown))
        print("Registered:", ", ".join(ENGINES))
        sys.exit(2)
    engines = [ENGINES[n] for n in names]

    if update:
        run_pptx(ENGINES["python"], update_golden=True)
        print("Golden layout updated. Re-run without --update-golden to verify.")
        return

    per_engine = {e.name: run_pptx(e) for e in engines}
    shared = [run_editor(), run_other_paths()]

    print(NL + "=" * 62)
    for name, res in per_engine.items():
        print(NL + "-- engine: %s --" % name)
        for ok, label in res.checks:
            print(f"  [{'PASS' if ok else 'FAIL'}] {label}")
        for sk in res.skips:
            print(f"  [SKIP] {sk}")
    print(NL + "-- engine-independent --")
    for res in shared:
        for ok, label in res.checks:
            print(f"  [{'PASS' if ok else 'FAIL'}] {label}")
        for sk in res.skips:
            print(f"  [SKIP] {sk}")

    # ---- per-engine summary table ----
    print(NL + "=" * 62)
    print("PDF -> PPTX, per engine")
    print(f"  {'engine':<10} {'ships':<7} {'passed':>8}  {'failed':>7}  verdict")
    print("  " + "-" * 56)
    blocking = False
    for e in engines:
        res = per_engine[e.name]
        ok = sum(1 for k, _ in res.checks if k)
        bad = sum(1 for k, _ in res.checks if not k)
        if not res.checks:
            verdict = "NOT RUN"
        elif bad == 0:
            verdict = "GREEN"
        elif e.ships:
            verdict = "RED (SHIPS: BLOCKING)"
            blocking = True
        else:
            verdict = "RED (quarantined, not linked)"
        print(f"  {e.name:<10} {'yes' if e.ships else 'no':<7} {ok:>8}  {bad:>7}  {verdict}")

    shared_ok = all(k for r in shared for k, _ in r.checks)
    if not shared_ok:
        blocking = True
    total = sum(len(r.checks) for r in list(per_engine.values()) + shared)
    passed = sum(1 for r in list(per_engine.values()) + shared for k, _ in r.checks if k)
    print("=" * 62)
    print(f"{passed}/{total} checks passed across {len(engines)} engine(s)")
    print("FIDELITY SUITE: " + ("GREEN - every shipping engine holds every invariant"
                                if not blocking else
                                "RED - a SHIPPING engine failed. Do not ship."))
    sys.exit(1 if blocking else 0)


if __name__ == "__main__":
    main()
