"""Region-level fallback verification for the browser engine.

The statement class that used to ship whole pages as "preserved as image":

  * a Type3-font statement (the AFP/print-stream generator class) must now
    convert fully native - pdf.js reports NaN metrics for Type3 fonts, and
    the NaN once poisoned every span bbox on the page;
  * a statement carrying a corrupt overprint stamp must ship ONLY the stamp
    region as image, keep the ruled table native and the rest editable, and
    say so honestly in the report note.

The bar, held here so it cannot regress: on this class at least 90% of the
source's text characters remain EDITABLE text in the produced PPTX, counted
by extracting text from the deck and comparing character counts against the
source page. Browser engine only: the desktop engine has no fallback path.

Synthetic only: every value is fabricated here. Never commit a real statement.
"""

import base64
import os
import random
import re
import sys
import zlib

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, os.path.join(HERE, "..", ".."))

import fitz
from pptx import Presentation

from tests.engines import BrowserEngine, EngineUnavailable, _Server

WORK = os.path.join(HERE, "..", "..", "_work", "region_fallback")
EDITABLE_BAR = 0.90

TYPE3_CHARS = "0123456789ABCDEFGHIJKLMNOPQRSTUVWXYZ.- %"


# --------------------------------------------------------------------------- #
# builders
# --------------------------------------------------------------------------- #
def _statement_body(page, rng):
    W = page.rect.width
    left, right = 40, W - 40
    cols = [left, 110, 330, 405, 480, right]
    page.insert_text((left, 60), "EXAMPLE BANK", fontname="hebo", fontsize=16)
    page.insert_text((left, 86), "A N Example - Account 00-00-00 12345678",
                     fontname="helv", fontsize=9)
    top, row_h, n_rows = 110, 20, 30
    bottom = top + row_h * (n_rows + 1)
    for r in range(n_rows + 2):
        yy = top + r * row_h
        page.draw_line(fitz.Point(left, yy), fitz.Point(right, yy), width=0.7)
    for x in cols:
        page.draw_line(fitz.Point(x, top), fitz.Point(x, bottom), width=0.7)
    page.draw_rect(fitz.Rect(left, top, right, top + row_h),
                   fill=(0.92, 0.92, 0.92), width=0)
    for ci, htext in enumerate(["Date", "Description", "Money out", "Money in", "Balance"]):
        page.insert_text((cols[ci] + 4, top + 13), htext, fontname="hebo", fontsize=8)
    balance = 4180.55
    for r in range(n_rows):
        yy = top + (r + 1) * row_h + 13
        amt = round(rng.uniform(3.5, 220.0), 2)
        balance -= amt
        page.insert_text((cols[0] + 4, yy), "%02d MAY" % (r % 28 + 1), fontname="helv", fontsize=8)
        page.insert_text((cols[1] + 4, yy), "FABRICATED PAYEE %02d" % (r % 12), fontname="helv", fontsize=8)
        page.insert_text((cols[2] + 4, yy), "%.2f" % amt, fontname="helv", fontsize=8)
        page.insert_text((cols[4] + 4, yy), "%.2f" % balance, fontname="helv", fontsize=8)
    page.insert_text((left, bottom + 24),
                     "Interest rate 1.20% AER variable. Fabricated data for testing only.",
                     fontname="helv", fontsize=7)


def make_stamped_statement(path):
    """Ruled statement plus a corrupt overprint stamp: eight DIFFERENT runs
    written at one origin (a generator writing its COPY stamp block without
    advancing the cursor). Only the stamp region may ship as image."""
    doc = fitz.open()
    page = doc.new_page(width=595, height=842)
    _statement_body(page, random.Random(5))
    for k in range(8):
        page.insert_text((430, 66), "DUPLICATE COPY %d" % k,
                         fontname="hebo", fontsize=10)
    doc.save(path)
    doc.close()


def make_type3_statement(path):
    """Hand-built statement whose text uses a Type3 font, the construct
    AFP/print-stream statement converters emit. Glyphs are crude filled
    shapes; the ToUnicode CMap keeps the text extractable."""
    rng = random.Random(4)
    W_GLYPH = 500

    def glyph_proc(ch):
        n = (ord(ch) * 7) % 300
        return ("%d 0 d0 50 0 %d %d re f" % (W_GLYPH, W_GLYPH - 100, 600 - n)).encode()

    def tounicode():
        pairs = "".join("<%02X> <%04X>\n" % (i + 33, ord(c))
                        for i, c in enumerate(TYPE3_CHARS))
        return ("/CIDInit /ProcSet findresource begin\n12 dict begin\nbegincmap\n"
                "/CMapName /T3U def\n/CMapType 2 def\n1 begincodespacerange\n"
                "<21> <%02X>\nendcodespacerange\n%d beginbfchar\n%sendbfchar\n"
                "endcmap\nCMapName currentdict /CMap defineresource pop\nend\nend"
                % (33 + len(TYPE3_CHARS) - 1, len(TYPE3_CHARS), pairs)).encode()

    def enc(s):
        return "".join("\\%03o" % (33 + TYPE3_CHARS.index(c))
                       for c in s.upper() if c in TYPE3_CHARS)

    W, H = 595, 842
    left, right = 40, W - 40
    cols = [left, 110, 330, 405, 480, right]
    top, row_h, n_rows = 110, 20, 30
    bottom = top + row_h * (n_rows + 1)

    c = [b"0.7 w"]
    for r in range(n_rows + 2):
        yy = H - (top + r * row_h)
        c.append(("%g %g m %g %g l S" % (left, yy, right, yy)).encode())
    for x in cols:
        c.append(("%g %g m %g %g l S" % (x, H - top, x, H - bottom)).encode())
    c.append(("0.92 g %g %g %g %g re f 0 g"
              % (left, H - top - row_h, right - left, row_h)).encode())

    def text(x, y, s, size):
        c.append(b"BT /F1 %d Tf %g %g Td (%s) Tj ET"
                 % (size, x, H - y, enc(s).encode()))

    text(left, 60, "EXAMPLE BANK", 16)
    text(right - 150, 60, "STATEMENT OF ACCOUNT", 9)
    text(left, 86, "A N EXAMPLE - ACCOUNT 00-00-00 12345678", 9)
    for ci, htext in enumerate(["DATE", "DESCRIPTION", "MONEY OUT", "MONEY IN", "BALANCE"]):
        text(cols[ci] + 4, top + 14, htext, 8)
    balance = 4180.55
    for r in range(n_rows):
        yy = top + (r + 1) * row_h + 14
        amt = round(rng.uniform(3.5, 220.0), 2)
        balance -= amt
        text(cols[0] + 4, yy, "%02d MAY" % (r % 28 + 1), 8)
        text(cols[1] + 4, yy, "FABRICATED PAYEE %02d" % (r % 12), 8)
        text(cols[2] + 4, yy, "%.2f" % amt, 8)
        text(cols[4] + 4, yy, "%.2f" % balance, 8)
    text(left, bottom + 24, "INTEREST RATE 1.20% AER VARIABLE. FABRICATED DATA.", 7)

    content = zlib.compress(b"\n".join(c))
    objs = {}
    onum = 10
    charproc_refs = []
    for i, ch in enumerate(TYPE3_CHARS):
        data = glyph_proc(ch)
        objs[onum] = b"<< /Length %d >>\nstream\n%s\nendstream" % (len(data), data)
        charproc_refs.append(onum)
        onum += 1
    cp = b"<< " + b" ".join(b"/c%02X %d 0 R" % (33 + i, num)
                            for i, num in enumerate(charproc_refs)) + b" >>"
    diffs = b"[ 33 " + b" ".join(b"/c%02X" % (33 + i)
                                 for i in range(len(TYPE3_CHARS))) + b" ]"
    widths = b"[ " + b" ".join(b"%d" % W_GLYPH for _ in TYPE3_CHARS) + b" ]"
    tu = tounicode()
    objs[5] = b"<< /Length %d >>\nstream\n%s\nendstream" % (len(tu), tu)
    objs[6] = (b"<< /Type /Font /Subtype /Type3 /FontBBox [0 -200 600 800] "
               b"/FontMatrix [0.001 0 0 0.001 0 0] /CharProcs " + cp +
               b" /Encoding << /Type /Encoding /Differences " + diffs + b" >>"
               b" /FirstChar 33 /LastChar %d /Widths " % (33 + len(TYPE3_CHARS) - 1) +
               widths + b" /ToUnicode 5 0 R >>")
    objs[4] = (b"<< /Length %d /Filter /FlateDecode >>\nstream\n" % len(content)
               + content + b"\nendstream")
    objs[3] = (b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 595 842] "
               b"/Resources << /Font << /F1 6 0 R >> >> /Contents 4 0 R >>")
    objs[2] = b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>"
    objs[1] = b"<< /Type /Catalog /Pages 2 0 R >>"

    out = bytearray(b"%PDF-1.4\n")
    offsets = {}
    for num in sorted(objs):
        offsets[num] = len(out)
        out += b"%d 0 obj\n" % num + objs[num] + b"\nendobj\n"
    xref_pos = len(out)
    maxnum = max(objs)
    out += b"xref\n0 %d\n" % (maxnum + 1)
    out += b"0000000000 65535 f \n"
    for num in range(1, maxnum + 1):
        if num in offsets:
            out += b"%010d 00000 n \n" % offsets[num]
        else:
            out += b"0000000000 65535 f \n"
    out += (b"trailer\n<< /Size %d /Root 1 0 R >>\nstartxref\n%d\n%%%%EOF\n"
            % (maxnum + 1, xref_pos))
    open(path, "wb").write(bytes(out))


# --------------------------------------------------------------------------- #
# conversion through the suite harness, report included
# --------------------------------------------------------------------------- #
def convert_with_report(engine, src_pdf, out_pptx):
    if not engine.available():
        raise EngineUnavailable("needs playwright and the benchpdf-site checkout")
    import threading
    from playwright.sync_api import sync_playwright

    srv = _Server(("127.0.0.1", 0), engine._handler())
    port = srv.server_address[1]
    threading.Thread(target=srv.serve_forever, daemon=True).start()
    try:
        with sync_playwright() as p:
            b = p.chromium.launch()
            pg = b.new_page()
            pg.goto(f"http://127.0.0.1:{port}/harness.html", wait_until="load")
            pg.wait_for_function("window.harnessReady === true", timeout=30000)
            data = list(open(src_pdf, "rb").read())
            r = pg.evaluate("bytes => window.convertToPptxWithReport(bytes)", data)
            b.close()
        with open(out_pptx, "wb") as f:
            f.write(base64.b64decode(r["pptx"]))
        return r["report"]
    finally:
        srv.shutdown()
        srv.server_close()


# --------------------------------------------------------------------------- #
def nchars(s):
    return len(re.sub(r"\s+", "", s))


def deck_editable_chars(pptx_path):
    prs = Presentation(pptx_path)
    total = 0
    for slide in prs.slides:
        for sh in slide.shapes:
            if sh.has_text_frame:
                total += nchars(sh.text_frame.text)
            elif sh.has_table:
                for row in sh.table.rows:
                    for cell in row.cells:
                        total += nchars(cell.text)
    return total


def native_tables(pptx_path):
    prs = Presentation(pptx_path)
    return sum(1 for sl in prs.slides for sh in sl.shapes if sh.has_table)


def src_chars(pdf_path):
    doc = fitz.open(pdf_path)
    n = sum(nchars(doc[i].get_text("text")) for i in range(doc.page_count))
    doc.close()
    return n


def main():
    os.makedirs(WORK, exist_ok=True)
    eng = BrowserEngine()
    checks = []

    def check(ok, label):
        checks.append((bool(ok), label))
        print("  [%s] %s" % ("PASS" if ok else "FAIL", label))

    # ---- Type3 statement: fully native, no fallback at all ----
    t3 = os.path.join(WORK, "type3_statement.pdf")
    make_type3_statement(t3)
    t3_out = os.path.join(WORK, "type3_statement.pptx")
    print("== Type3 statement (AFP/print-stream class) ==")
    rep = convert_with_report(eng, t3, t3_out)
    modes = [p["mode"] for p in rep["pages"]]
    check(modes == ["native"], f"converts native, no fallback (modes: {modes})")
    check(not rep.get("notes"), f"no fallback notes (notes: {rep.get('notes')})")
    check(native_tables(t3_out) >= 1, "ruled table arrives native")
    ratio = deck_editable_chars(t3_out) / max(src_chars(t3), 1)
    check(ratio >= 0.98, f"text fully editable ({ratio:.1%})")

    # ---- stamped statement: region fallback, body stays editable ----
    st = os.path.join(WORK, "stamped_statement.pdf")
    make_stamped_statement(st)
    st_out = os.path.join(WORK, "stamped_statement.pptx")
    print("== stamped statement (corrupt overprint region) ==")
    rep = convert_with_report(eng, st, st_out)
    modes = [p["mode"] for p in rep["pages"]]
    check(modes == ["region-fallback"],
          f"region fallback, not page fallback (modes: {modes})")
    notes = rep.get("notes") or []
    check(len(notes) == 1 and re.fullmatch(
        r"\d+ regions? on page 1 preserved as image", notes[0] or ""),
        f"honest region note (notes: {notes})")
    check(native_tables(st_out) >= 1, "ruled table still native")
    ratio = deck_editable_chars(st_out) / max(src_chars(st), 1)
    check(ratio >= EDITABLE_BAR,
          f"at least {EDITABLE_BAR:.0%} of characters stay editable ({ratio:.1%})")

    bad = sum(1 for ok, _ in checks if not ok)
    print("\n%d/%d region-fallback checks passed" % (len(checks) - bad, len(checks)))
    return 1 if bad else 0


if __name__ == "__main__":
    sys.exit(main())
