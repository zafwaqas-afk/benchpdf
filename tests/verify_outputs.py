"""
Programmatic verification of real conversion outputs produced by verify_hub.py
(which drives the actual browser + server). Opens each output with the
appropriate library and checks it is structurally valid.
"""
import glob
import os
import sys
import zipfile

import fitz
from docx import Document
from pptx import Presentation

HERE = os.path.dirname(os.path.abspath(__file__))
WORK = os.path.join(HERE, "..", "_work")

results = []


def check(name, cond, detail=""):
    results.append((name, cond, detail))
    print(f"  [{'PASS' if cond else 'FAIL'}] {name}{(' - ' + detail) if detail else ''}")


def latest_file(pattern):
    matches = sorted(glob.glob(pattern), key=os.path.getmtime)
    return matches[-1] if matches else None


def main():
    # PDFs produced by Office COM export (Word/Excel/PowerPoint -> PDF)
    for label, name_glob in [
        ("Word -> PDF output", "vendor_checklist.pdf"),
        ("Excel -> PDF output", "budget.pdf"),
        ("PowerPoint -> PDF output", "hub_test_deck.pdf"),
    ]:
        p = latest_file(os.path.join(WORK, "**", name_glob))
        if not p:
            check(label, False, "no output file found")
            continue
        try:
            d = fitz.open(p)
            check(label, d.page_count > 0, f"{os.path.basename(p)}: {d.page_count} page(s)")
            d.close()
        except Exception as exc:
            check(label, False, str(exc))

    # Images -> merged PDF (should be 3 pages: jpg+png+heic)
    p = latest_file(os.path.join(WORK, "**", "merged.pdf"))
    if p:
        d = fitz.open(p)
        check("Images -> merged PDF page count", d.page_count == 3, f"{d.page_count} pages (expect 3)")
        d.close()
    else:
        check("Images -> merged PDF", False, "no output found")

    # Web page -> PDF
    p = latest_file(os.path.join(WORK, "**", "*example*.pdf")) or latest_file(os.path.join(WORK, "**", "*.pdf"))
    web_hits = glob.glob(os.path.join(WORK, "**", "example*.pdf"), recursive=True)
    if web_hits:
        p = sorted(web_hits, key=os.path.getmtime)[-1]
        d = fitz.open(p)
        txt = d[0].get_text("text")
        check("Web page -> PDF has real page content", "Example Domain" in txt, txt[:60].replace("\n", " "))
        d.close()
    else:
        check("Web page -> PDF", False, "no output found")

    # PDF -> PowerPoint (existing engine)
    p = latest_file(os.path.join(WORK, "**", "tables_charts.pptx"))
    if p:
        prs = Presentation(p)
        check("PDF -> PowerPoint slide count", len(prs.slides._sldIdLst) >= 1,
              f"{len(prs.slides._sldIdLst)} slide(s)")
    else:
        check("PDF -> PowerPoint", False, "no output found")

    # PDF -> Word (round trip via Word's reflow import)
    p = latest_file(os.path.join(WORK, "**", "text_report.docx"))
    if p:
        try:
            doc = Document(p)
            paras = [x.text for x in doc.paragraphs if x.text.strip()]
            check("PDF -> Word opens and has real text", len(paras) > 0, f"{len(paras)} paragraph(s)")
        except Exception as exc:
            check("PDF -> Word opens with python-docx", False, str(exc))
    else:
        check("PDF -> Word", False, "no output found (see PDF->Word timing note)")

    # PDF -> Images (zip of PNGs)
    p = latest_file(os.path.join(WORK, "**", "tables_charts.zip"))
    if p:
        with zipfile.ZipFile(p) as z:
            names = z.namelist()
        check("PDF -> Images zip contains page images", len(names) >= 1, f"{len(names)} file(s): {names}")
    else:
        check("PDF -> Images", False, "no output found")

    # PDF -> Text
    p = latest_file(os.path.join(WORK, "**", "text_report.txt"))
    if p:
        txt = open(p, encoding="utf-8").read()
        check("PDF -> Text has real content", len(txt) > 100, f"{len(txt)} chars")
    else:
        check("PDF -> Text", False, "no output found")

    print(f"\n{sum(1 for _, ok, _ in results if ok)}/{len(results)} output checks passed")
    sys.exit(0 if all(ok for _, ok, _ in results) else 1)


if __name__ == "__main__":
    main()
