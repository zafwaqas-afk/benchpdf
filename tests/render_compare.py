"""
Visual fidelity + editability verification (real render, not code review).

  * Renders PPTX slides 1-3 to PNG by driving installed PowerPoint via COM.
  * Renders the same PDF pages to PNG at the same pixel size.
  * Composes side-by-side comparison images (PDF | PPTX).
  * Editability: appends a sentence to the first body paragraph on slide 1,
    saves, re-renders, and shows that the text re-wraps inside its box.
"""
import os
import sys

import fitz
import win32com.client
import pythoncom
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import MSO_AUTO_SIZE

HERE = os.path.dirname(os.path.abspath(__file__))
# Pass a PDF path as the first argument; defaults to a bundled synthetic sample.
SRC = sys.argv[1] if len(sys.argv) > 1 else os.path.join(HERE, "input", "tables_charts.pdf")
_STEM = os.path.splitext(os.path.basename(SRC))[0]
OUT = os.path.join(HERE, "output", _STEM + "_reworked.pptx")
EDIT = os.path.join(HERE, "output", _STEM + "_reworked_edited.pptx")
CMP = os.path.join(HERE, "compare")
os.makedirs(CMP, exist_ok=True)

W = 820  # target render width in px


def render_pdf_page(page, width):
    z = width / page.rect.width
    pix = page.get_pixmap(matrix=fitz.Matrix(z, z), alpha=False)
    return Image.frombytes("RGB", (pix.width, pix.height), pix.samples)


def export_slides(pptx_path, indices, width):
    """Export given 1-based slide indices to PNG via PowerPoint COM."""
    pythoncom.CoInitialize()
    pp = win32com.client.Dispatch("PowerPoint.Application")
    out = {}
    try:
        pres = pp.Presentations.Open(os.path.abspath(pptx_path), ReadOnly=True,
                                     Untitled=False, WithWindow=False)
        sw = pres.PageSetup.SlideWidth
        sh = pres.PageSetup.SlideHeight
        height = int(round(width * sh / sw))
        for n in indices:
            p = os.path.join(CMP, f"_slide{n}.png")
            pres.Slides(n).Export(p, "PNG", width, height)
            out[n] = p
        pres.Close()
    finally:
        pp.Quit()
        pythoncom.CoUninitialize()
    return out


def label(img, text):
    strip = Image.new("RGB", (img.width, 26), (28, 26, 22))
    d = ImageDraw.Draw(strip)
    d.text((8, 7), text, fill=(230, 226, 216))
    return strip


def compose(pdf_img, ppt_img, title, path):
    h = max(pdf_img.height, ppt_img.height)
    for im in (pdf_img, ppt_img):
        if im.height < h:
            bg = Image.new("RGB", (im.width, h), (255, 255, 255))
            bg.paste(im, (0, 0))
    gap = 16
    lab_pdf = label(pdf_img, "PDF (source)")
    lab_ppt = label(ppt_img, "PPTX (converted, rendered by PowerPoint)")
    colw = pdf_img.width
    canvas = Image.new("RGB", (pdf_img.width + ppt_img.width + gap, h + 26 + 30),
                       (245, 243, 238))
    d = ImageDraw.Draw(canvas)
    d.text((8, 8), title, fill=(30, 27, 20))
    canvas.paste(lab_pdf, (0, 30))
    canvas.paste(pdf_img, (0, 56))
    canvas.paste(lab_ppt, (pdf_img.width + gap, 30))
    canvas.paste(ppt_img, (pdf_img.width + gap, 56))
    canvas.save(path)
    return path


def ink_iou(a, b):
    """Rough structural overlap: IoU of dark-ink pixels at low resolution."""
    import numpy as np
    s = (160, 200)
    ga = np.asarray(a.convert("L").resize(s)) < 128
    gb = np.asarray(b.convert("L").resize(s)) < 128
    inter = (ga & gb).sum()
    union = (ga | gb).sum()
    return inter / union if union else 1.0


def main():
    doc = fitz.open(SRC)
    slides = export_slides(OUT, [1, 2, 3], W)
    print("=== visual fidelity (pages 1-3) ===")
    for n in (1, 2, 3):
        pdf_img = render_pdf_page(doc[n - 1], W)
        ppt_img = Image.open(slides[n]).convert("RGB")
        path = compose(pdf_img, ppt_img, f"Page {n}", os.path.join(CMP, f"compare_page{n}.png"))
        print(f"  page {n}: ink-overlap IoU = {ink_iou(pdf_img, ppt_img)*100:.1f}%  -> {path}")

    # ---- editability check ----
    print("\n=== editability (append + wrap) ===")
    prs = Presentation(OUT)
    slide1 = prs.slides[0]
    target = None
    for sh in slide1.shapes:
        if sh.has_text_frame and not sh.has_table:
            size = max([r.font.size.pt for p in sh.text_frame.paragraphs
                        for r in p.runs if r.font.size], default=0)
            txt = sh.text_frame.text.strip()
            if 6 <= size <= 13 and len(txt) > 40:
                target = sh
                break
    assert target is not None, "no body text box found on slide 1"
    para = target.text_frame.paragraphs[0]
    before_lines_hint = target.text_frame.text
    run = para.add_run()
    run.text = (" Appended sentence: confirming that the body text re-flows and wraps "
                "within its box across multiple lines instead of overflowing one frozen line.")
    if para.runs and para.runs[0].font.size:
        run.font.size = para.runs[0].font.size
        run.font.name = para.runs[0].font.name
    word_wrap = target.text_frame.word_wrap
    autosize = target.text_frame.auto_size
    prs.save(EDIT)
    print(f"  target box word_wrap={word_wrap}, auto_size={autosize} "
          f"(None = no autofit, so appended text wraps)")
    ok_edit = (word_wrap is True) and (autosize == MSO_AUTO_SIZE.NONE or autosize is None)

    edited = export_slides(EDIT, [1], W)
    img = Image.open(edited[1]).convert("RGB")
    demo = os.path.join(CMP, "edit_demo.png")
    lab = label(img, "Slide 1 after appending a sentence — text wraps inside the box")
    canvas = Image.new("RGB", (img.width, img.height + 26), (245, 243, 238))
    canvas.paste(lab, (0, 0)); canvas.paste(img, (0, 26)); canvas.save(demo)
    print(f"  wrote {demo}")
    print(f"  [{'PASS' if ok_edit else 'FAIL'}] appended text will wrap (word_wrap on, autofit off)")

    doc.close()
    print("\nDONE")


if __name__ == "__main__":
    main()
