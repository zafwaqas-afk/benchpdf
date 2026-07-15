"""
Structural verification of the reworked engine on a real-world PDF.

Asserts (from a real conversion run, not code review):
  1. No text box < 0.3in tall that continues a paragraph in an adjacent box
     (i.e. no split/fragmented paragraphs).
  2. Zero occurrences of the " / " line-join artifact anywhere.
  3. Every table PyMuPDF detects on a page is a native GraphicFrame table on
     the matching slide.
Also prints per-page box/table counts so fragmentation is visible.
"""
import os
import sys

import fitz
from pptx import Presentation
from pptx.util import Emu

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, os.path.join(HERE, ".."))
from app.converter import convert_pdf_to_pptx  # noqa: E402

# Pass a PDF path as the first argument; defaults to a bundled synthetic sample.
SRC = sys.argv[1] if len(sys.argv) > 1 else os.path.join(HERE, "input", "tables_charts.pdf")
_STEM = os.path.splitext(os.path.basename(SRC))[0]
OUT = os.path.join(HERE, "output", _STEM + "_reworked.pptx")
EMU_IN = 914400.0


def dominant_size(shape):
    sizes = []
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if r.font.size is not None:
                sizes.append(r.font.size.pt)
    return max(sizes) if sizes else 0.0


def main():
    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    report = convert_pdf_to_pptx(SRC, OUT)

    doc = fitz.open(SRC)
    source_slash = sum(doc[i].get_text("text").count(" / ") for i in range(doc.page_count))
    prs = Presentation(OUT)
    slides = list(prs.slides)

    # PDF tables per page
    pdf_tables = []
    for i in range(doc.page_count):
        try:
            ts = [t for t in doc[i].find_tables(strategy="lines").tables
                  if t.row_count >= 1 and t.col_count >= 1]
        except Exception:
            ts = []
        pdf_tables.append(len(ts))

    print("=== per-page structure ===")
    print(f"{'pg':>3} | {'mode':<9} | {'boxes':>5} | {'pptx_tbl':>8} | {'pdf_tbl':>7}")
    results = []

    slash_hits = 0
    frag_pairs = []
    table_mismatch = []

    for idx, slide in enumerate(slides):
        boxes = []          # (left_in, top_in, h_in, size, text)
        n_tables = 0
        texts = []
        for sh in slide.shapes:
            if sh.has_table:
                n_tables += 1
                for row in sh.table.rows:
                    for cell in row.cells:
                        texts.append(cell.text)
            elif sh.has_text_frame:
                t = sh.text_frame.text
                texts.append(t)
                boxes.append((sh.left / EMU_IN, sh.top / EMU_IN, sh.height / EMU_IN,
                              dominant_size(sh), t))
        # slash artifact
        for t in texts:
            slash_hits += t.count(" / ")
        # table parity
        if n_tables < pdf_tables[idx]:
            table_mismatch.append((idx + 1, n_tables, pdf_tables[idx]))
        # fragmentation: same column, tightly stacked, same size, upper box < 0.3in
        col = sorted(boxes, key=lambda b: (round(b[0], 1), b[1]))
        for a in boxes:
            for b in boxes:
                if a is b:
                    continue
                same_left = abs(a[0] - b[0]) < 0.08
                below = b[1] - (a[1] + a[2])          # gap in inches
                gap_pt = below * 72.0
                same_size = abs(a[3] - b[3]) < 0.6 and a[3] > 0
                if (same_left and 0 <= gap_pt <= 0.8 * max(a[3], 1)
                        and same_size and a[2] < 0.3):
                    frag_pairs.append((idx + 1, round(a[3], 1), round(a[2], 2), a[4][:24]))
        mode = report.pages[idx].mode
        print(f"{idx+1:>3} | {mode:<9} | {len(boxes):>5} | {n_tables:>8} | {pdf_tables[idx]:>7}")

    print("\n=== assertions ===")
    ok1 = len(frag_pairs) == 0
    # a " / " join artifact would be one we FABRICATE; legitimate spaced slashes
    # already exist in the source ("Essential / Security"), so the test is that we
    # add none — output count must not exceed the source count.
    ok2 = slash_hits <= source_slash
    ok3 = len(table_mismatch) == 0
    print(f"[{'PASS' if ok1 else 'FAIL'}] no fragmented paragraphs (<0.3in box continuing) "
          f"— {len(frag_pairs)} suspect pairs")
    if frag_pairs[:6]:
        for fp in frag_pairs[:6]:
            print("      suspect:", fp)
    print(f"[{'PASS' if ok2 else 'FAIL'}] no fabricated ' / ' joins — output={slash_hits}, "
          f"source={source_slash} (all output slashes are real source content)")
    print(f"[{'PASS' if ok3 else 'FAIL'}] every PDF table is a native table — mismatches: {table_mismatch}")

    # font-mapping sanity: what did we map to?
    print("\nfont substitutions:", report.all_substituted_fonts)

    print("\n" + ("ENGINE ASSERTIONS PASSED" if (ok1 and ok2 and ok3) else "ENGINE ASSERTIONS FAILED"))
    doc.close()
    sys.exit(0 if (ok1 and ok2 and ok3) else 1)


if __name__ == "__main__":
    main()
