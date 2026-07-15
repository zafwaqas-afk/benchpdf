"""
Convert the three test PDFs and programmatically verify the spec requirements:

  * slide count == PDF page count
  * on native-mode pages, text is present as real text (not baked into images)
  * text-box positions are within 2% of the source PDF coordinates
  * round-trip token overlap >= 95% on the text-heavy sample
  * no picture-only "empty deck" for text PDFs

Prints a per-page conversion report for each file.
"""

import os
import re
import sys

import fitz
from pptx import Presentation
from pptx.util import Emu

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, os.path.join(HERE, ".."))
from app.converter import convert_pdf_to_pptx  # noqa: E402

INP = os.path.join(HERE, "input")
OUTP = os.path.join(HERE, "output")
os.makedirs(OUTP, exist_ok=True)

EMU_PER_PT = 12700


def tokenize(text):
    return [t for t in re.findall(r"[A-Za-z0-9]+", text.lower()) if t]


def token_overlap(a, b):
    from collections import Counter
    ca, cb = Counter(tokenize(a)), Counter(tokenize(b))
    if not ca:
        return 1.0 if not cb else 0.0
    inter = sum((ca & cb).values())
    return inter / sum(ca.values())


def print_report(report):
    print(f"\n=== Conversion report: {os.path.basename(report.source_pdf)} ===")
    print(f"pages: {report.page_count}")
    hdr = f"{'page':>4} | {'mode':<10} | {'boxes':>5} | {'imgs':>4} | {'vpaths':>6} | fonts substituted"
    print(hdr)
    print("-" * len(hdr))
    for p in report.pages:
        subs = ", ".join(p.substituted_fonts) if p.substituted_fonts else "-"
        note = f"   [{p.note}]" if p.note else ""
        print(f"{p.page_number:>4} | {p.mode:<10} | {p.text_boxes:>5} | {p.images:>4} | "
              f"{p.vector_paths:>6} | {subs}{note}")
    for w in report.warnings:
        print("  WARNING:", w)


def pptx_slide_text(slide):
    parts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            parts.append(shape.text_frame.text)
    return "\n".join(parts)


def verify_file(name, expect_modes=None):
    src = os.path.join(INP, name)
    out = os.path.join(OUTP, name.replace(".pdf", ".pptx"))
    report = convert_pdf_to_pptx(src, out)
    print_report(report)

    results = []

    def check(cond, label):
        results.append((cond, label))
        print(f"  [{'PASS' if cond else 'FAIL'}] {label}")

    doc = fitz.open(src)
    prs = Presentation(out)

    # 1. slide count == page count
    check(len(prs.slides.__iter__.__self__._sldIdLst) == doc.page_count
          if False else len(prs.slides._sldIdLst) == doc.page_count,
          f"slide count ({len(prs.slides._sldIdLst)}) == page count ({doc.page_count})")

    # 2. native pages have real text
    slides = list(prs.slides)
    for pr, slide in zip(report.pages, slides):
        if pr.mode in ("native", "hybrid"):
            txt = pptx_slide_text(slide).strip()
            check(len(txt) > 0, f"page {pr.page_number} ({pr.mode}): editable text present")

    # 3. position accuracy within 2% on the first native page
    #    compare first text block bbox in PDF vs first textbox in slide
    slide_w_pt = prs.slide_width / EMU_PER_PT
    slide_h_pt = prs.slide_height / EMU_PER_PT
    checked_pos = False
    for idx, (pr, slide) in enumerate(zip(report.pages, slides)):
        if pr.mode != "native":
            continue
        page = doc[idx]
        blocks = [b for b in page.get_text("dict")["blocks"]
                  if b.get("type", 0) == 0 and any(
                      sp.get("text", "").strip()
                      for ln in b.get("lines", []) for sp in ln.get("spans", []))]
        if not blocks:
            continue
        b0 = blocks[0]["bbox"]
        # first textbox shape on the slide
        tb = None
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                tb = shape
                break
        if tb is None:
            continue
        got_left = tb.left / EMU_PER_PT
        got_top = tb.top / EMU_PER_PT
        dx = abs(got_left - b0[0]) / slide_w_pt
        dy = abs(got_top - b0[1]) / slide_h_pt
        check(dx <= 0.02 and dy <= 0.02,
              f"page {pr.page_number}: first text box within 2% "
              f"(dx={dx*100:.2f}%, dy={dy*100:.2f}%)")
        checked_pos = True
        break
    if not checked_pos and expect_modes and "native" in expect_modes:
        check(False, "expected a native page to check position accuracy")

    # expected mode presence
    if expect_modes:
        modes_seen = {p.mode for p in report.pages}
        for m in expect_modes:
            check(m in modes_seen, f"expected a '{m}' page present (saw {sorted(modes_seen)})")

    doc.close()
    return report, results


def main():
    all_ok = True

    _, r1 = verify_file("text_report.pdf", expect_modes=["native"])
    _, r2 = verify_file("slide_deck.pdf", expect_modes=["native"])
    _, r3 = verify_file("tables_charts.pdf", expect_modes=["hybrid"])

    # round-trip token overlap on the text-heavy sample
    print("\n=== Round-trip token overlap (text_report.pdf) ===")
    src = os.path.join(INP, "text_report.pdf")
    out = os.path.join(OUTP, "text_report.pptx")
    doc = fitz.open(src)
    pdf_text = "\n".join(doc[i].get_text("text") for i in range(doc.page_count))
    doc.close()
    prs = Presentation(out)
    pptx_text = "\n".join(pptx_slide_text(s) for s in prs.slides)
    overlap = token_overlap(pdf_text, pptx_text)
    ok = overlap >= 0.95
    print(f"  token overlap = {overlap*100:.2f}%  [{'PASS' if ok else 'FAIL'}] (>=95%)")

    for _, res in (("", r1), ("", r2), ("", r3)):
        for cond, _label in res:
            all_ok = all_ok and cond
    all_ok = all_ok and ok

    print("\n" + ("ALL CHECKS PASSED" if all_ok else "SOME CHECKS FAILED"))
    sys.exit(0 if all_ok else 1)


if __name__ == "__main__":
    main()
