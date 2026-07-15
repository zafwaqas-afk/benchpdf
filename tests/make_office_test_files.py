"""Generate real Office test files for the conversion hub matrix: a .docx, an
.xlsx, a .pptx, a few images (jpg/png/heic), plus a password-protected .docx
and a corrupt "PDF" for the error-path tests."""
import os
import shutil

from docx import Document
from docx.shared import Pt
import xlsxwriter
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt
from PIL import Image, ImageDraw

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, "input")
os.makedirs(OUT, exist_ok=True)


def make_docx(path):
    d = Document()
    d.add_heading("Vendor Onboarding Checklist", level=1)
    d.add_paragraph(
        "This document tracks the steps required to onboard a new third-party "
        "vendor, from initial due diligence through to contract signature. "
        "Each step below should be completed in order.")
    d.add_heading("Due diligence", level=2)
    for item in ["Confirm company registration and ownership structure",
                 "Request references from two existing clients",
                 "Run a sanctions and adverse-media screening"]:
        d.add_paragraph(item, style="List Bullet")
    d.add_heading("Contracting", level=2)
    table = d.add_table(rows=3, cols=2)
    table.style = "Light Grid Accent 1"
    rows = [("Step", "Owner"), ("Draft MSA", "Legal"), ("Countersign", "Procurement")]
    for r, (a, b) in enumerate(rows):
        table.cell(r, 0).text = a
        table.cell(r, 1).text = b
    d.add_paragraph("End of checklist.")
    d.save(path)
    print("wrote", path)


def make_password_docx(path, password="secret123"):
    """Password-protect via Word COM SaveAs2 (python-docx can't set passwords)."""
    import win32com.client
    word = win32com.client.DispatchEx("Word.Application")
    word.Visible = False
    word.DisplayAlerts = 0
    try:
        doc = word.Documents.Add()
        doc.Content.Text = "This document is password protected for testing purposes."
        # Late-bound COM doesn't reliably honour SaveAs2's Password= kwarg;
        # setting the Password property directly is the documented VBA-equivalent
        # way to set the open-password before saving.
        doc.Password = password
        doc.SaveAs2(os.path.abspath(path), FileFormat=12)
        doc.Close(False)
    finally:
        word.Quit()
    print("wrote", path, "(password:", password, ")")


def make_xlsx(path):
    wb = xlsxwriter.Workbook(path)
    ws = wb.add_worksheet("Budget")
    bold = wb.add_format({"bold": True, "bg_color": "#1f4e79", "font_color": "white"})
    headers = ["Category", "Q1", "Q2", "Q3", "Q4"]
    for c, h in enumerate(headers):
        ws.write(0, c, h, bold)
    rows = [("Salaries", 120, 125, 130, 128),
            ("Marketing", 40, 55, 60, 45),
            ("Infrastructure", 30, 32, 35, 38)]
    for r, row in enumerate(rows, start=1):
        for c, v in enumerate(row):
            ws.write(r, c, v)
    ws.set_column(0, 0, 18)
    wb.close()
    print("wrote", path)


def make_pptx(path):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.6), Inches(8.5), Inches(1.2))
    tb.text_frame.text = "Office Conversion Hub — Test Deck"
    tb.text_frame.paragraphs[0].runs[0].font.size = PptPt(32)
    tb2 = slide.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(8.5), Inches(1.0))
    tb2.text_frame.text = "Used to verify PPTX -> PDF via PowerPoint COM automation."
    prs.save(path)
    print("wrote", path)


def make_images(out_dir):
    made = []
    im1 = Image.new("RGB", (900, 600), (235, 225, 200))
    d = ImageDraw.Draw(im1)
    d.rectangle([40, 40, 860, 560], outline=(90, 70, 30), width=6)
    d.text((80, 80), "Test image A (JPG)", fill=(40, 30, 10))
    p1 = os.path.join(out_dir, "photo_a.jpg")
    im1.save(p1, "JPEG", quality=90)
    made.append(p1)

    im2 = Image.new("RGB", (700, 900), (200, 220, 230))
    d2 = ImageDraw.Draw(im2)
    d2.ellipse([60, 60, 640, 840], outline=(20, 60, 90), width=8)
    d2.text((100, 100), "Test image B (PNG)", fill=(10, 40, 60))
    p2 = os.path.join(out_dir, "photo_b.png")
    im2.save(p2, "PNG")
    made.append(p2)

    # HEIC via pillow-heif
    import pillow_heif
    pillow_heif.register_heif_opener()
    im3 = Image.new("RGB", (800, 600), (220, 200, 210))
    d3 = ImageDraw.Draw(im3)
    d3.rectangle([30, 30, 770, 570], outline=(90, 30, 50), width=6)
    d3.text((70, 70), "Test image C (HEIC)", fill=(60, 10, 20))
    p3 = os.path.join(out_dir, "photo_c.heic")
    im3.save(p3, format="HEIF", quality=85)
    made.append(p3)
    return made


def make_corrupt_pdf(path):
    with open(path, "wb") as f:
        f.write(b"%PDF-1.4\nThis is not a real PDF body, just garbage bytes to trigger a parse failure.\n%%EOF")
    print("wrote", path, "(deliberately corrupt)")


if __name__ == "__main__":
    make_docx(os.path.join(OUT, "vendor_checklist.docx"))
    make_xlsx(os.path.join(OUT, "budget.xlsx"))
    make_pptx(os.path.join(OUT, "hub_test_deck.pptx"))
    make_images(OUT)
    make_corrupt_pdf(os.path.join(OUT, "corrupt.pdf"))
    make_password_docx(os.path.join(OUT, "protected.docx"))
    print("done")
