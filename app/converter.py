"""
PDF -> editable PPTX conversion engine.

This module owns ONE thing: the PPTX *placement policy*. All PDF extraction
(line clustering into logical paragraphs, document-wide font mapping, table-cell
fill sampling, geometry) lives in app/extraction.py and is shared with the
editor. Here we take those extracted blocks and lay them down as positioned
PowerPoint shapes and native tables. Keep placement here; keep extraction there.

One PDF page => one slide. The engine works structurally, not by screenshotting:

  * TABLES   are detected from the PDF's vector ruling lines and rebuilt as NATIVE
             PowerPoint tables (rows/columns/merges, per-cell fill sampled from the
             PDF, per-cell text with real fonts). Table regions are excluded from
             any background render so nothing is drawn twice.
  * TEXT     spans are clustered into LOGICAL BLOCKS (paragraphs / headings) that
             become native, editable text boxes with word-wrap on. Wrapped lines are
             joined with a single space -- never a separator character.
  * IMAGES   embedded rasters are placed as pictures at their true position/size.
  * HYBRID   only genuinely non-tabular vector art falls back to an image render
             with editable text on top.
"""

from __future__ import annotations

import io
from dataclasses import dataclass, field
from typing import Callable, Optional

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn

# ---- shared extraction layer (do not re-implement these here) ----
from app.extraction import (
    EMU_PER_PT, SAMPLE_DPI, FLAG_ITALIC, FLAG_BOLD,
    FontMapper, _int_color_to_rgb, _center, _point_in, _sample_fill,
    _collect_lines, _line_alignment, _attach_markers, _cluster_lines,
    _split_paragraphs, _infer_aligned_tables, _inherit_glyph_colors,
    _glyph_outline_drawings, _text_width_pt,
)

# --------------------------------------------------------------------------- #
# PPTX-placement-only constants
# --------------------------------------------------------------------------- #
HYBRID_DPI = 200            # render resolution for hybrid backgrounds
NONTABLE_ART_COVERAGE = 0.03   # >3% of page area in non-table drawings -> hybrid
# "No Style, Table Grid" - thin borders, no theme fills (we set fills per cell).
TABLE_GRID_STYLE = "{5940675A-B579-460E-94D1-54222C63F5DA}"
# "No Style, No Grid" - for INFERRED (unruled) tables: the source drew no
# ruling lines, so the output must not invent them.
TABLE_NO_GRID_STYLE = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"
PROSE_CELL_CHARS = 300      # a cell holding this much text holds a paragraph
MIN_TABULAR_CELLS = 2       # one populated cell states no relationship
PROSE_MAX_TEXT_CELLS = 3    # ... and neither does a paragraph plus a label
NOWRAP_MAX_WORDS = 5        # blocks this short keep their source line breaks
WORD_FIT_SAFETY = 1.1       # substituted fonts can run a little wider
WRAP_SLACK_PT = 2.0         # breathing room at the box edge (source points)
TRACK_MAX_EM = 0.08         # never squeeze or open more than this per char


# --------------------------------------------------------------------------- #
# Report data structures
# --------------------------------------------------------------------------- #
@dataclass
class PageReport:
    page_number: int
    mode: str
    text_boxes: int = 0
    tables: int = 0
    images: int = 0
    vector_paths: int = 0
    substituted_fonts: list = field(default_factory=list)
    note: str = ""


@dataclass
class ConversionReport:
    source_pdf: str
    output_pptx: str
    page_count: int = 0
    pages: list = field(default_factory=list)
    scanned_warning: bool = False
    warnings: list = field(default_factory=list)

    @property
    def all_substituted_fonts(self) -> list:
        seen = {}
        for p in self.pages:
            for s in p.substituted_fonts:
                seen[s] = True
        return list(seen.keys())


# --------------------------------------------------------------------------- #
# PPTX placement: text blocks
# --------------------------------------------------------------------------- #
def _line_tracking(ln, fonts: FontMapper, src_width: float) -> float:
    """Points per character to add (or remove) so this line occupies src_width.

    A wrapping paragraph is re-broken by PowerPoint, and the metric-compatible
    substitute is never exactly the source font, so it wraps at different words
    and the column drifts. Measured against the corpus, the PDF's own line
    breaks are worth about 0.025 of real median over reflowed ones, and the
    drift is what compounds into blocks overflowing their box.

    Rather than give up reflow - which is what keeps the text editable, a
    paragraph the reader can retype and have re-wrap - each source line gets
    the tracking that restores its original width, so PowerPoint's greedy wrap
    breaks it where the source broke it. Clamped hard, because a bad
    measurement must degrade to slightly-wrong spacing, never to unreadable
    text; unmeasurable text gets no tracking at all.
    """
    if src_width <= 0:
        return 0.0
    measured = 0.0
    chars = 0
    max_size = 0.0
    for sp in ln["spans"]:
        text = sp.get("text") or ""
        if not text:
            continue
        flags = int(sp.get("flags", 0))
        size = float(sp.get("size", 10.0)) or 10.0
        measured += _text_width_pt(text, fonts.map(sp.get("font", ""), flags), size,
                                   bool(flags & FLAG_BOLD), bool(flags & FLAG_ITALIC))
        chars += len(text)
        max_size = max(max_size, size)
    if not measured or chars < 2:
        return 0.0
    limit = TRACK_MAX_EM * (max_size or 10.0)
    return max(-limit, min(limit, (src_width - measured) / chars))


def _emit_paragraph(text_frame, para_lines, first, alignment, fonts: FontMapper,
                    space_before: float = 0.0, track_lines: bool = False):
    """Write one paragraph's lines as a single wrapped paragraph (space-joined)."""
    p = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
    if alignment is not None:
        p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(0)
    prev_text = ""
    for li, ln in enumerate(para_lines):
        if li > 0 and prev_text and not prev_text[-1].isspace():
            # join wrapped lines with a single space - never a separator character
            first_span = ln["spans"][0]
            sep = p.add_run()
            sep.text = " "
            _style_run(sep, first_span, fonts)
        # Only lines that actually WRAPPED are tracked, which is every line of
        # the paragraph but its last. A line that ended because the paragraph
        # ended - a list item, a closing line - never broke, gains nothing from
        # being restored to its width, and only picks up distorted letter
        # spacing. Tracking every line cost boe_mpr_2025_08 p3 0.108 and
        # w3c_svg10_2001 0.021 of mean while helping the prose pages.
        wrapped = li < len(para_lines) - 1
        track = (_line_tracking(ln, fonts, ln["x1"] - ln["x0"])
                 if (track_lines and wrapped) else 0.0)
        for sp in ln["spans"]:
            run = p.add_run()
            run.text = sp["text"]
            _style_run(run, sp, fonts)
            if track:
                run.font._rPr.set("spc", str(int(round(track * 100))))
            prev_text = sp["text"]
    return p


def _style_run(run, span, fonts: FontMapper):
    f = run.font
    f.name = fonts.map(span.get("font", ""), int(span.get("flags", 0)))
    size = float(span.get("size", 10.0))
    if size > 0:
        f.size = Pt(size)   # preserve the exact point size
    flags = int(span.get("flags", 0))
    f.bold = bool(flags & FLAG_BOLD)
    f.italic = bool(flags & FLAG_ITALIC)
    try:
        f.color.rgb = _int_color_to_rgb(span.get("color", 0))
    except Exception:
        pass


def _longest_word_width(cluster) -> float:
    """Chars-proportional estimate of the widest single word in the block:
    enough to guarantee the box can hold it, which stops mid-word wrapping."""
    max_w = 0.0
    for ln in cluster:
        text = "".join(s["text"] for s in ln["spans"])
        char_w = (ln["x1"] - ln["x0"]) / max(len(text), 1)
        for wd in text.split():
            max_w = max(max_w, len(wd) * char_w)
    return max_w


def _source_leading(cluster) -> float:
    """The source's own line pitch, in points, or 0 when there isn't one to read.

    PowerPoint otherwise sets its own leading from the font, which is tighter
    than most documents': the W3C working draft leads 14pt text at 18.7pt, so a
    25-item list ended 130pt above where it should and the whole page read as
    vertically compressed. Bounded, because a block whose lines are not evenly
    pitched (a heading over a paragraph) has no single leading to impose.
    """
    if len(cluster) < 2:
        return 0.0
    deltas = sorted(d for d in (cluster[i]["y0"] - cluster[i - 1]["y0"]
                                for i in range(1, len(cluster))) if d > 0)
    if len(deltas) < 2:
        return 0.0
    med = deltas[len(deltas) // 2]
    size = max((l["size"] for l in cluster), default=1.0) or 1.0
    if med < size * 0.9 or med > size * 2.5:
        return 0.0
    return med


def _paragraph_gaps(paras, lead, box_height) -> list:
    """The air the source leaves above each paragraph, in points.

    A block used to emit its paragraphs hard against each other, because
    space_before was pinned to zero. Documents separate paragraphs with a
    blank line: the SEC chairman's letter (the corpus's worst page) leaves
    about half a line above each one, and without it the whole column runs
    together and ends short. What the source leaves is whatever exceeds the
    block's own leading, so it is only readable once that leading is known.

    The gaps are clamped to the room the source actually left for them. The box
    is sized to the source line span, which already includes those gaps, but
    PowerPoint stacks the lines at the full leading AND then adds the gaps on
    top - so the raw gaps overshoot the box and push the last line out the
    bottom. On w3c_svg10_2001 p2 that dropped the final line onto the heading
    below: 35 lines * 15.2pt = 532pt of pitch left only 112pt for gaps in a
    644pt box, but the raw gaps summed 127. Scale them to fit the 112.
    """
    gaps = [0.0] * len(paras)
    if not lead:
        return gaps
    for i in range(1, len(paras)):
        extra = paras[i][0]["y0"] - paras[i - 1][-1]["y0"] - lead
        gaps[i] = extra if extra > 0.5 else 0.0
    n_lines = sum(len(p) for p in paras)
    avail = box_height - n_lines * lead
    total = sum(gaps)
    if total > avail and total > 0:
        scale = max(avail, 0.0) / total
        gaps = [g * scale for g in gaps]
    return gaps


def _paragraph_indents(paras, block_x0) -> list:
    """Where each paragraph sits inside its block, in points from the block's
    left edge: (left_indent, first_line_indent).

    A list is one text box of paragraphs, and every paragraph used to start at
    the box's left edge. On the W3C working draft that flattened two nesting
    levels onto one margin and threw away the hanging indent that puts a
    wrapped line clear of its own bullet. Both are in the source geometry: the
    first line's x0 gives the outdent, the continuation lines' give the margin.
    """
    out = []
    for para in paras:
        first_x = para[0]["x0"] - block_x0
        rest = para[1:]
        # left_indent applies to every line; first_line_indent shifts only the
        # first, and is negative for a hanging indent, positive for a
        # first-line one.
        cont_x = (min(l["x0"] for l in rest) - block_x0) if rest else first_x
        mar_l = max(cont_x, 0.0)
        out.append((mar_l, first_x - mar_l))
    return out


def _add_text_block(slide, cluster, scale, off_x, off_y, fonts: FontMapper,
                    page_w: float = 0.0):
    x0 = min(c["x0"] for c in cluster)
    y0 = min(c["y0"] for c in cluster)
    x1 = max(c["x1"] for c in cluster)
    y1 = max(c["y1"] for c in cluster)

    words = sum(len("".join(s["text"] for s in ln["spans"]).split()) for ln in cluster)
    # Short header blocks ("END OF DAY / ACCOUNT BALANCE") must never re-wrap:
    # keep the source's own line breaks and switch wrapping off entirely.
    no_wrap_short = len(cluster) > 1 and words <= NOWRAP_MAX_WORDS
    wrap = len(cluster) > 1 and not no_wrap_short

    # a wrapping box must at minimum fit its longest word, or PowerPoint
    # breaks mid-word; cap at the page's right edge
    w_pt = max(x1 - x0, 1.0)
    if wrap:
        w_pt = max(w_pt, WORD_FIT_SAFETY * _longest_word_width(cluster))
        # A hair of breathing room. The box is pinned to the widest source
        # line's exact extent, so a full line whose substitute width lands
        # within a tenth of a point of that extent tips over PowerPoint's own
        # metrics - which are not PyMuPDF's - and wraps one line early. Too
        # small to drag a word across a line; just enough to clear the metric
        # divergence at the box edge.
        w_pt += WRAP_SLACK_PT
        if page_w:
            w_pt = min(w_pt, max(page_w - x0, x1 - x0))

    left = Emu(int((off_x + x0 * scale) * EMU_PER_PT))
    top = Emu(int((off_y + y0 * scale) * EMU_PER_PT))
    width = Emu(max(int(w_pt * scale * EMU_PER_PT), EMU_PER_PT))
    height = Emu(max(int((y1 - y0) * scale * EMU_PER_PT), EMU_PER_PT // 2))

    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE          # never let PPT change the font size
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    lead = _source_leading(cluster)
    alignment = _line_alignment(cluster, x0, x1)
    paras, indents = [], None
    if no_wrap_short:
        # one paragraph per source line: the layout is the source's, verbatim
        for li, ln in enumerate(cluster):
            paras.append(_emit_paragraph(tf, [ln], li == 0, alignment, fonts))
    else:
        split = _split_paragraphs(cluster)
        gaps = _paragraph_gaps(split, lead, y1 - y0)
        for pi, para in enumerate(split):
            paras.append(_emit_paragraph(tf, para, pi == 0, alignment, fonts,
                                         gaps[pi] * scale, track_lines=wrap))
        indents = _paragraph_indents(split, x0)
    if lead:
        for p in paras:
            p.line_spacing = Pt(lead * scale)
    if indents:
        # python-pptx exposes no left/first-line indent on a paragraph, so the
        # attributes go straight onto a:pPr - the same two attributes the
        # browser engine stamps into its OOXML.
        for p, (mar_l, first) in zip(paras, indents):
            if abs(mar_l) > 0.5 or abs(first) > 0.5:
                pPr = p._p.get_or_add_pPr()
                pPr.set("marL", str(int(round(mar_l * scale * EMU_PER_PT))))
                pPr.set("indent", str(int(round(first * scale * EMU_PER_PT))))
    return tb


# --------------------------------------------------------------------------- #
# PPTX placement: native tables
# --------------------------------------------------------------------------- #
def _set_table_grid_style(table, style_id=TABLE_GRID_STYLE):
    tblPr = table._tbl.tblPr
    if tblPr is None:
        return
    for child in list(tblPr):
        if child.tag == qn("a:tableStyleId"):
            tblPr.remove(child)
    el = tblPr.makeelement(qn("a:tableStyleId"), {})
    el.text = style_id
    tblPr.append(el)


def _lines_in(bbox, lines):
    out = []
    for ln in lines:
        cx, cy = _center(ln["bbox"])
        if _point_in(bbox, cx, cy):
            out.append(ln)
    return out


def _table_cells_tabular(table, lines) -> bool:
    """A bordered callout is not a table.

    Judged on cell CONTENT, so it applies to every table however it was
    detected. A govuk guidance page shipped a ruled 2x2 holding one
    1,600-character paragraph in one cell, two cells empty and the page number
    in the last: a box drawn around prose, whose text reflows inside a native
    table into a column one cell wide. A table earns its cells by using them:
    two of them must carry text, and a cell-sized paragraph is only tabular
    when enough other cells answer it.
    """
    cells = [c for row in table.rows for c in row.cells if c is not None]
    if not cells:
        return False
    lens = [sum(len(s["text"].strip()) for ln in _lines_in(c, lines) for s in ln["spans"])
            for c in cells]
    text_cells = sum(1 for n in lens if n > 0)
    if text_cells < MIN_TABULAR_CELLS:
        return False
    if text_cells <= PROSE_MAX_TEXT_CELLS and max(lens) >= PROSE_CELL_CHARS:
        return False
    return True


def _fill_cell(cell, para_lines, fonts: FontMapper, fill_rgb, is_dark):
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(*fill_rgb)
    cell.margin_left = Pt(4)
    cell.margin_right = Pt(4)
    cell.margin_top = Pt(1)
    cell.margin_bottom = Pt(1)
    cell.vertical_anchor = MSO_ANCHOR.TOP

    tf = cell.text_frame
    tf.word_wrap = True
    if not para_lines:
        # An empty cell still carries a paragraph, and without an explicit size
        # PowerPoint reserves line height for its 18pt default, silently
        # inflating every sparse row until dense tables overflow the slide.
        # (Found by the phase-3 render comparison against the browser engine.)
        tf.paragraphs[0].font.size = Pt(6)
        return
    para_lines = sorted(para_lines, key=lambda l: (round(l["y0"], 1), l["x0"]))
    paras = _split_paragraphs(para_lines)
    for pi, para in enumerate(paras):
        _emit_paragraph(tf, para, pi == 0, None, fonts)


def _add_table(slide, table, page_lines, pix, z, scale, off_x, off_y, fonts: FontMapper):
    nrows, ncols = table.row_count, table.col_count
    if nrows < 1 or ncols < 1:
        return False
    bx0, by0, bx1, by1 = table.bbox

    left = Emu(int((off_x + bx0 * scale) * EMU_PER_PT))
    top = Emu(int((off_y + by0 * scale) * EMU_PER_PT))
    width = Emu(int((bx1 - bx0) * scale * EMU_PER_PT))
    height = Emu(int((by1 - by0) * scale * EMU_PER_PT))

    gf = slide.shapes.add_table(nrows, ncols, left, top, width, height)
    table_obj = gf.table
    table_obj.first_row = False
    table_obj.horz_banding = False
    _set_table_grid_style(table_obj,
                          TABLE_NO_GRID_STYLE if getattr(table, "inferred", False)
                          else TABLE_GRID_STYLE)

    row0 = table.rows[0].cells
    for ci in range(ncols):
        cb = row0[ci]
        if cb is not None:
            table_obj.columns[ci].width = Emu(max(int((cb[2] - cb[0]) * scale * EMU_PER_PT), EMU_PER_PT))
    for ri in range(nrows):
        rcells = [c for c in table.rows[ri].cells if c is not None]
        if rcells:
            rh = max(c[3] for c in rcells) - min(c[1] for c in rcells)
            table_obj.rows[ri].height = Emu(max(int(rh * scale * EMU_PER_PT), EMU_PER_PT // 3))

    for ri in range(nrows):
        for ci in range(ncols):
            cb = table.rows[ri].cells[ci]
            cell = table_obj.cell(ri, ci)
            if cb is None:
                cell.fill.background()
                cell.text_frame.paragraphs[0].font.size = Pt(6)
                continue
            fill = _sample_fill(pix, cb, z)
            is_dark = (fill[0] + fill[1] + fill[2]) < 384
            cell_lines = _lines_in(cb, page_lines)
            _fill_cell(cell, cell_lines, fonts, fill, is_dark)
    return True


# --------------------------------------------------------------------------- #
# Hybrid background (only for genuine non-table vector art)
# --------------------------------------------------------------------------- #
def _nontable_art_ratio(page, drawings, table_bboxes) -> float:
    page_area = abs(page.rect.width * page.rect.height) or 1.0
    area = 0.0
    for d in drawings:
        r = d.get("rect")
        if r is None:
            continue
        w, h = r.width, r.height
        if w <= 2 or h <= 2:
            continue
        if w > 0.92 * page.rect.width and h > 0.92 * page.rect.height:
            continue
        cx, cy = (r.x0 + r.x1) / 2, (r.y0 + r.y1) / 2
        if any(_point_in(tb, cx, cy) for tb in table_bboxes):
            continue
        area += w * h
    return area / page_area


def _render_hybrid_bg(src_doc, page_index, table_bboxes, glyph_rects=()) -> bytes:
    zoom = HYBRID_DPI / 72.0
    tmp = fitz.open()
    tmp.insert_pdf(src_doc, from_page=page_index, to_page=page_index)
    pg = tmp[0]
    try:
        pg.add_redact_annot(pg.rect, text=None)
        pg.apply_redactions(images=fitz.PDF_REDACT_IMAGE_NONE,
                            graphics=fitz.PDF_REDACT_LINE_ART_NONE)
        # Second pass: text painted as filled glyph-OUTLINE paths (invisible-
        # text-layer statements). Redaction only removes text objects, so the
        # painted word would survive as pixels while its invisible twin ships
        # editable on top: ghost-doubled text. Cover each matched outline blob
        # and remove line art fully inside it; background fills extend beyond
        # the blob's bbox and survive.
        if glyph_rects:
            for r in glyph_rects:
                pg.add_redact_annot(fitz.Rect(*r) + (-1, -1, 1, 1), text=None)
            pg.apply_redactions(images=fitz.PDF_REDACT_IMAGE_NONE,
                                graphics=fitz.PDF_REDACT_LINE_ART_REMOVE_IF_COVERED)
    except Exception:
        tmp.close()
        tmp = fitz.open()
        tmp.insert_pdf(src_doc, from_page=page_index, to_page=page_index)
        pg = tmp[0]
    for tb in table_bboxes:
        try:
            pg.draw_rect(fitz.Rect(*tb), color=(1, 1, 1), fill=(1, 1, 1))
        except Exception:
            pass
    pix = pg.get_pixmap(matrix=fitz.Matrix(zoom, zoom), alpha=False)
    data = pix.tobytes("png")
    tmp.close()
    return data


def _lock_picture(picture):
    try:
        pic = picture._element
        cNvPicPr = pic.find(qn("p:nvPicPr")).find(qn("p:cNvPicPr"))
        locks = cNvPicPr.find(qn("a:picLocks"))
        if locks is None:
            locks = cNvPicPr.makeelement(qn("a:picLocks"), {})
            cNvPicPr.append(locks)
        for attr in ("noMove", "noResize", "noSelect", "noChangeAspect"):
            locks.set(attr, "1")
    except Exception:
        pass


# --------------------------------------------------------------------------- #
# Main entry point
# --------------------------------------------------------------------------- #
def convert_pdf_to_pptx(
    pdf_path: str,
    pptx_path: str,
    progress_callback: Optional[Callable[[int, int, str], None]] = None,
) -> ConversionReport:
    report = ConversionReport(source_pdf=pdf_path, output_pptx=pptx_path)
    fonts = FontMapper()

    doc = fitz.open(pdf_path)
    report.page_count = doc.page_count
    if doc.page_count == 0:
        report.warnings.append("The PDF has no pages.")
        doc.close()
        return report

    first = doc[0]
    slide_w_pt, slide_h_pt = first.rect.width, first.rect.height
    prs = Presentation()
    prs.slide_width = Emu(int(slide_w_pt * EMU_PER_PT))
    prs.slide_height = Emu(int(slide_h_pt * EMU_PER_PT))
    blank = prs.slide_layouts[6]

    pages_with_text = 0

    for i in range(doc.page_count):
        page = doc[i]
        if progress_callback:
            progress_callback(i, doc.page_count, f"Converting page {i + 1} of {doc.page_count}")
        slide = prs.slides.add_slide(blank)

        pw, ph = page.rect.width, page.rect.height
        scale = min(slide_w_pt / pw, slide_h_pt / ph) if pw and ph else 1.0
        off_x = (slide_w_pt - pw * scale) / 2.0
        off_y = (slide_h_pt - ph * scale) / 2.0

        text_dict = page.get_text("dict")
        raw_text = page.get_text("text").strip()
        try:
            drawings = page.get_drawings()
        except Exception:
            drawings = []

        # ---- scanned / image-only ----
        if not raw_text:
            pix = page.get_pixmap(matrix=fitz.Matrix(HYBRID_DPI / 72.0, HYBRID_DPI / 72.0), alpha=False)
            bg = slide.shapes.add_picture(
                io.BytesIO(pix.tobytes("png")),
                Emu(int(off_x * EMU_PER_PT)), Emu(int(off_y * EMU_PER_PT)),
                Emu(int(pw * scale * EMU_PER_PT)), Emu(int(ph * scale * EMU_PER_PT)))
            _lock_picture(bg)
            report.pages.append(PageReport(page_number=i + 1, mode="image-only",
                                           vector_paths=len(drawings),
                                           note="No extractable text (scanned/image-only page)."))
            report.scanned_warning = True
            continue

        pages_with_text += 1

        all_lines = _collect_lines(text_dict)
        # run colours: text painted as glyph-outline paths carries its real
        # colour on the paths, not on the (often invisible) text layer
        _inherit_glyph_colors(all_lines, drawings, pw, ph)

        # ---- tables ----
        try:
            found = page.find_tables(strategy="lines")
            detected = [t for t in found.tables if t.row_count >= 1 and t.col_count >= 1]
        except Exception:
            detected = []
        # A detected grid must contain text to be promoted to a native table:
        # a decorative squares mark is graphics, not an empty 1xN table.
        def _has_cell_text(t):
            return any(_point_in(t.bbox, *_center(ln["bbox"])) for ln in all_lines)
        ruled = [t for t in detected if _has_cell_text(t)]
        demoted_grids = len(detected) - len(ruled)
        # Unruled tables (statement ledgers without ruling lines) are
        # recovered from column alignment and emitted native.
        inferred = _infer_aligned_tables(all_lines, [t.bbox for t in ruled])
        # A table whose cells hold prose rather than tabular content is
        # furniture: its border stays in the background layer and its text
        # flows as ordinary text boxes.
        tables = [t for t in ruled + inferred if _table_cells_tabular(t, all_lines)]
        table_bboxes = [t.bbox for t in tables]

        loose_lines = [ln for ln in all_lines
                       if not any(_point_in(tb, *_center(ln["bbox"])) for tb in table_bboxes)]

        # ---- hybrid decision (non-table vector art only) ----
        # A demoted text-less grid forces the hybrid background so the
        # decoration still ships, as pixels in the background layer.
        art_ratio = _nontable_art_ratio(page, drawings, table_bboxes)
        use_hybrid = art_ratio > NONTABLE_ART_COVERAGE or demoted_grids > 0
        pr = PageReport(page_number=i + 1, mode="hybrid" if use_hybrid else "native",
                        vector_paths=len(drawings))

        if use_hybrid:
            pr.note = f"non-table vector art ({art_ratio*100:.0f}% of page)"
            if demoted_grids:
                pr.note += f"; {demoted_grids} decorative grid(s) kept as graphics"
            glyph_rects = [r for r, _ in
                           _glyph_outline_drawings(drawings, all_lines, pw * ph)]
            png = _render_hybrid_bg(doc, i, table_bboxes, glyph_rects)
            bg = slide.shapes.add_picture(
                io.BytesIO(png),
                Emu(int(off_x * EMU_PER_PT)), Emu(int(off_y * EMU_PER_PT)),
                Emu(int(pw * scale * EMU_PER_PT)), Emu(int(ph * scale * EMU_PER_PT)))
            _lock_picture(bg)
        else:
            for img in page.get_images(full=True):
                xref = img[0]
                try:
                    rects = page.get_image_rects(xref)
                except Exception:
                    rects = []
                for rect in rects:
                    cx, cy = (rect.x0 + rect.x1) / 2, (rect.y0 + rect.y1) / 2
                    if any(_point_in(tb, cx, cy) for tb in table_bboxes):
                        continue
                    try:
                        base = doc.extract_image(xref)
                        slide.shapes.add_picture(
                            io.BytesIO(base["image"]),
                            Emu(int((off_x + rect.x0 * scale) * EMU_PER_PT)),
                            Emu(int((off_y + rect.y0 * scale) * EMU_PER_PT)),
                            Emu(int(rect.width * scale * EMU_PER_PT)),
                            Emu(int(rect.height * scale * EMU_PER_PT)))
                        pr.images += 1
                    except Exception:
                        pass

        # ---- native tables ----
        if tables:
            pix = page.get_pixmap(matrix=fitz.Matrix(SAMPLE_DPI / 72.0, SAMPLE_DPI / 72.0), alpha=False)
            z = SAMPLE_DPI / 72.0
            for t in tables:
                if _add_table(slide, t, all_lines, pix, z, scale, off_x, off_y, fonts):
                    pr.tables += 1

        # ---- loose text as logical blocks ----
        loose_lines = _attach_markers(loose_lines)
        for cluster in _cluster_lines(loose_lines):
            _add_text_block(slide, cluster, scale, off_x, off_y, fonts, pw)
            pr.text_boxes += 1

        pr.substituted_fonts = [f"{k} -> {v}" for k, v in sorted(fonts.substitutions.items())]
        report.pages.append(pr)

    if progress_callback:
        progress_callback(doc.page_count, doc.page_count, "Saving presentation")

    if report.scanned_warning and pages_with_text == 0:
        report.warnings.append(
            "This PDF appears to be scanned or image-only (no selectable text). Pages were "
            "placed as images so nothing is lost, but there is no editable text to extract. "
            "OCR is not part of this tool.")
    elif report.scanned_warning:
        report.warnings.append("Some pages had no extractable text and were placed as images.")

    all_subs = [f"{k} -> {v}" for k, v in sorted(fonts.substitutions.items())]
    for p in report.pages:
        if p.mode != "image-only":
            p.substituted_fonts = all_subs

    prs.save(pptx_path)
    doc.close()
    return report
